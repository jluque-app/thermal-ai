# app_improved.py
import sys
import os
# Hack for Render: Ensure current directory is in sys.path so sibling imports work
# regardless of start command (uvicorn backend.app vs cd backend && uvicorn app)
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

import threading
import uuid
from typing import Optional, Dict, Any, List, Tuple
import io
import os
from pathlib import Path
import traceback
import inspect
from datetime import datetime
import json
import time
import hmac
import hashlib

print("DEBUG: Force Deploy - Route Fix Verification 001")

import sys
import os
import json
import base64
import io
import time
import logging
import uuid
import shutil
import requests
from datetime import datetime, timedelta, timezone
from typing import Optional, Dict, List, Any, Union
from pathlib import Path

from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Header, Depends, Request, Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, EmailStr

# -----------------------------
# Optional ML Imports (Handle Vercel Limits)
# -----------------------------
try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    cv2 = None
    CV2_AVAILABLE = False

try:
    import torch
    import torchvision.transforms as transforms
    from segment_anything import sam_model_registry, SamAutomaticMaskGenerator
    ML_AVAILABLE = True
except ImportError:
    print("WARNING: ML libraries (torch/segment_anything) not found. Running in LITE mode.")
    ML_AVAILABLE = False
    torch = None
    transforms = None
    sam_model_registry = None
    SamAutomaticMaskGenerator = None

    SamAutomaticMaskGenerator = None

# -----------------------------
# Core Image Libs (Made Optional for Vercel "Ultra-Lite")
# -----------------------------
try:
    import numpy as np
    from PIL import Image, ImageDraw, ImageFont
    IMAGE_LIBS_AVAILABLE = True
except ImportError:
    print("WARNING: Numpy/PIL not found. Image processing disabled.")
    IMAGE_LIBS_AVAILABLE = False
    np = None
    Image = None
    ImageDraw = None
    ImageFont = None

import openai
import stripe
import openai
import stripe

# -----------------------------
# Reporting Libs (Made Optional for Vercel "Extreme-Lite")
# -----------------------------
import openai
import stripe

# -----------------------------
# Reporting Libs (Made Optional for Vercel "Extreme-Lite")
# -----------------------------
try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches, Pt
    REPORTING_LIBS_AVAILABLE = True
except ImportError:
    print("WARNING: Reporting libraries (reportlab/docx/pptx) not found. Report generation disabled.")
    REPORTING_LIBS_AVAILABLE = False
    canvas = None
    A4 = None
    Document = None
    Presentation = None
    Inches = None
    Pt = None

from pdf_endpoint import router as pdf_router
from docx_endpoint import router as docx_router
# NOTE: we do NOT include ppt_endpoint router to avoid route conflicts with inline /v1/report/ppt
# from ppt_endpoint import router as ppt_router

from expert_chat_endpoint import router as expert_chat_router
from expert_lead_endpoint import router as expert_lead_router

from segmentation_utils import SegmentationModel
from thermal_core_improved import (
    detect_hotspot_mask,
    overlay_mask_on_rgb,
    encode_image_to_base64_png,
    calculate_external_heat_transfer_coefficient,
    instantaneous_loss_proxy_watts,
    annualize_proxy_kwh,
    infer_u_value,
    annual_kwh_saved_u_method,
    annual_total_loss_u_method,
    compute_multi_year_costs,
    MIN_RELIABLE_DELTA_T_C,
)

# Audited multi-modal registration engine (registration_engine.py, v2)
try:
    from registration_engine import MultiModalImageAligner, ENGINE_VERSION as REGISTRATION_ENGINE_VERSION
    HAVE_REGISTRATION_ENGINE = True
except Exception as _reg_exc:  # pragma: no cover
    print("WARNING: registration_engine unavailable:", repr(_reg_exc))
    MultiModalImageAligner = None
    REGISTRATION_ENGINE_VERSION = None
    HAVE_REGISTRATION_ENGINE = False

BACKEND_VERSION = "v2026.08.28-audit-integration"
DEFAULT_KAPPA_CALIBRATION = 0.05   # dimensionless screening calibration constant
DEFAULT_WIND_SPEED_M_S = 1.5

from climate_data_improved import get_outdoor_temperature_c, degree_hours_below_base
from report_template_improved import build_gamma_payload
from report_builder import build_report


# -----------------------------
# App (CREATE FIRST!)
# -----------------------------
app = FastAPI(
    title="ThermalAI Backend (Improved)",
    version="0.2.9",  # bumped
    docs_url="/docs",
    openapi_url="/openapi.json",
    redoc_url=None,
)

@app.get("/__ppt_template_inspect")
def ppt_template_inspect():
    """
    Returns the list of shapes in the PPT template with their identifiers,
    so we can see how IMG_RGB / IMG_OVERLAY placeholders are labeled.
    """
    try:
        from pptx import Presentation

        template_path = _template_pptx_path()
        prs = Presentation(template_path)

        out = []
        for si, slide in enumerate(prs.slides, start=1):
            for sh in slide.shapes:
                # Alt text descr (if exists)
                descr = ""
                cname = ""
                try:
                    cNvPr = sh.element.xpath(".//p:cNvPr")[0]
                    descr = (cNvPr.get("descr") or "").strip()
                    cname = (cNvPr.get("name") or "").strip()
                except Exception:
                    pass

                text = ""
                try:
                    if getattr(sh, "has_text_frame", False) and sh.text_frame:
                        text = (sh.text_frame.text or "").strip()
                except Exception:
                    pass

                out.append({
                    "slide": si,
                    "shape_name": getattr(sh, "name", ""),
                    "cNvPr_name": cname,
                    "alt_descr": descr,
                    "has_text": bool(text),
                    "text": text[:80],
                    "shape_type": int(getattr(sh, "shape_type", 0)),
                })

        return {"template": template_path, "count": len(out), "shapes": out}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e), "traceback": traceback.format_exc()})


# -----------------------------
# Routers (AFTER app exists)
# -----------------------------
app.include_router(pdf_router)
app.include_router(docx_router)
# app.include_router(ppt_router)  # DISABLED (we use inline PPT endpoint below)
app.include_router(expert_chat_router)
app.include_router(expert_lead_router)

# -----------------------------
# CORS (ONLY ONCE)
# -----------------------------
ALLOWED_ORIGINS = [
    "https://thermal-insights-1d99c858.base44.app",
    "https://app.thermalai.eu",
    "https://thermalai.eu",
    "http://localhost:3000",
    "http://localhost:5173",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

MODEL_PATH = Path("models/Model.pth")
SEG_MODEL: Optional[SegmentationModel] = None


# -----------------------------
# Health/debug endpoints
# -----------------------------
# -----------------------------
# Static Files (Frontend)
# -----------------------------
# Verify if frontend/dist exists (Docker build puts it at /app/frontend/dist)
# The script runs from /app/backend, so we look up one level
frontend_dist = Path(__file__).resolve().parent.parent / "frontend" / "dist"

print(f"DEBUG: app_improved.py is at {Path(__file__).resolve()}")
print(f"DEBUG: Looking for frontend at {frontend_dist}")
print(f"DEBUG: Exists? {frontend_dist.exists()}")
if frontend_dist.exists():
    print(f"DEBUG: Contents of dist: {[f.name for f in frontend_dist.iterdir()]}")

if frontend_dist.exists():
    app.mount("/assets", StaticFiles(directory=str(frontend_dist / "assets")), name="assets")
    
    # Catch-all route moved to end of file to avoid shadowing API routes
    pass



@app.get("/health")
def health() -> Dict[str, Any]:
    return {"status": "ok", "model_loaded": SEG_MODEL is not None}


@app.get("/__whoami")
def whoami() -> Dict[str, Any]:
    return {"app": "app_improved.py", "version": "0.2.9"}


# -----------------------------
# Model loading (non-blocking)
# -----------------------------
def ensure_model_downloaded() -> None:
    if MODEL_PATH.exists():
        return

    url = os.getenv("MODEL_GDRIVE_URL", "").strip()
    if not url:
        raise RuntimeError("MODEL_GDRIVE_URL is not set (Render Environment Variable).")

    MODEL_PATH.parent.mkdir(parents=True, exist_ok=True)

    import torch
    print(f"Downloading Model from {url}...")
    if "drive.google.com" in url:
        import gdown
        gdown.download(url=url, output=str(MODEL_PATH), quiet=False, fuzzy=True)
    else:
        # Direct download (e.g. from Meta FB servers)
        torch.hub.download_url_to_file(url, str(MODEL_PATH))
    
    print("Model download complete:", str(MODEL_PATH))


@app.on_event("startup")
def startup_event() -> None:
    threading.Thread(target=_load_model_background, daemon=True).start()
    print("Startup: model loading kicked off in background.")


def _load_model_background() -> None:
    global SEG_MODEL
    try:
        ensure_model_downloaded()
        SEG_MODEL = SegmentationModel(model_path=str(MODEL_PATH))
        print("Segmentation model loaded.")
    except Exception as e:
        print("WARN: loading segmentation model failed, falling back to mock:", repr(e))
        try:
             # Fallback to mock if real model fails (e.g. missing file/URL)
             from segmentation_utils import MockSegmentationModel
             SEG_MODEL = MockSegmentationModel()
             print("MockSegmentationModel loaded (Fallback).")
        except Exception as e2:
             print("CRITICAL: Failed to load even MockSegmentationModel:", repr(e2))
             # Even if this fails, we create a dummy class so usage doesn't crash
             class DummyModel:
                 def to(self, d): return self
                 def generate(self, img): return []
                 def predict_masks(self, img):
                     # Return a fake result compatible with analyze()
                     import numpy as np
                     w, h = img.size
                     # Create simple masks (all wall, small window)
                     wall = np.ones((h, w), dtype=bool)
                     window = np.zeros((h, w), dtype=bool)
                     # Center window
                     cw, ch = w//2, h//2
                     window[ch-20:ch+20, cw-20:cw+20] = True
                     door = np.zeros((h, w), dtype=bool)
                     counts = {
                         "wall_pixels": int(wall.sum()),
                         "window_pixels": int(window.sum()),
                         "door_pixels": int(door.sum()),
                         "total_pixels": w*h
                     }
                     class FakeSegResult:
                         def __init__(self):
                             self.wall_mask = wall
                             self.window_mask = window
                             self.door_mask = door
                             self.counts = counts
                             self.indexed = np.zeros((h,w), dtype=np.uint8)
                     return FakeSegResult()
             SEG_MODEL = DummyModel()


# -----------------------------
# Helpers
# -----------------------------
def _safe_float(v: Optional[str]) -> Optional[float]:
    if v is None or v == "":
        return None
    try:
        return float(v)
    except Exception:
        return None


def _safe_bool(v: Optional[str], default: bool = False) -> bool:
    if v is None:
        return default
    return str(v).strip().lower() in ("true", "1", "yes", "y", "on")


def _resize_max(img: Image.Image, max_side: int = 1024) -> Image.Image:
    w, h = img.size
    m = max(w, h)
    if m <= max_side:
        return img
    scale = max_side / float(m)
    new_w = max(1, int(round(w * scale)))
    new_h = max(1, int(round(h * scale)))
    return img.resize((new_w, new_h))

# --------------------------------------------------
# Thermal → RGB registration helpers (cv2-based)
# --------------------------------------------------

def _pil_to_bgr(img: Image.Image) -> np.ndarray:
    """Convert RGB PIL → BGR numpy for OpenCV."""
    return cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)


def _bgr_to_pil(arr: np.ndarray) -> Image.Image:
    """Convert BGR numpy → RGB PIL."""
    return Image.fromarray(cv2.cvtColor(arr, cv2.COLOR_BGR2RGB))


def register_thermal_to_rgb(vis_img: Image.Image, thr_img: Image.Image):
    """
    Align thermal -> RGB with the audited MultiModalImageAligner (registration_engine.py).

    Returns:
        aligned_thermal (PIL.Image), valid_support_mask (np.ndarray bool, RGB frame),
        registration_meta (dict) with stage outcome, quality flags and diagnostics.
    """
    if HAVE_REGISTRATION_ENGINE and cv2 is not None:
        aligner = MultiModalImageAligner()
        res = aligner.register(vis_img, thr_img)
        md = res.metadata or {}
        score = md.get("alignment_score_final")
        if res.is_reliable and (score is None or score >= 0.30):
            quality_label = "high"
        elif res.is_reliable:
            quality_label = "medium"
        elif md.get("prior_source") in ("identity_same_dims", "exif_f35") or str(md.get("prior_source", "")).startswith("camera_table"):
            quality_label = "prior_only"
        else:
            quality_label = "poor"
        meta = {
            "engine_version": md.get("engine_version"),
            "used": bool(res.is_reliable),
            "method": res.method_used,
            "confidence": res.confidence_score,
            "reliable": bool(res.is_reliable),
            "quality_label": quality_label,
            "inliers": int(res.inlier_count),
            "matrix": res.homography_matrix.tolist() if res.homography_matrix is not None else None,
            "prior_source": md.get("prior_source"),
            "prior_scale": md.get("prior_scale"),
            "final_scale": md.get("final_scale"),
            "alignment_score_prior": md.get("alignment_score_prior"),
            "alignment_score_final": md.get("alignment_score_final"),
            "alignment_gain": md.get("alignment_gain"),
            "support_coverage": md.get("support_coverage"),
            "ecc_completed_levels": md.get("ecc_completed_levels"),
            "ecc_diagnostics": md.get("ecc_diagnostics"),
            "rejections": md.get("rejections"),
            "rgb_model": md.get("rgb_model"),
            "thermal_model": md.get("thermal_model"),
            "target_registration_error": None,   # only measurable against ground-truth landmarks
        }
        return res.aligned_thermal_pil, res.valid_support_mask, meta

    # Engine unavailable: plain resize with an honest "poor" flag.
    resized = thr_img.resize(vis_img.size, Image.Resampling.BILINEAR)
    valid_mask = np.ones((vis_img.size[1], vis_img.size[0]), dtype=bool)
    meta = {
        "engine_version": None, "used": False, "method": "resize_only", "confidence": None,
        "reliable": False, "quality_label": "poor", "inliers": 0, "matrix": None,
        "prior_source": None, "alignment_score_prior": None, "alignment_score_final": None,
        "alignment_gain": None, "support_coverage": 1.0, "ecc_completed_levels": 0,
        "ecc_diagnostics": [], "rejections": ["registration engine not available"],
        "target_registration_error": None,
    }
    return resized, valid_mask, meta



def _template_pptx_path() -> str:
    return str((Path(__file__).resolve().parent / "templates" / "ThermalAI.pptx").resolve())


def _get_ppt_builder():
    """
    ppt_report_builder has not been stable across versions.
    This resolver supports either:
      - build_reports(...)
      - build_report(...)
    """
    import ppt_report_builder as prb
    if hasattr(prb, "build_reports"):
        return prb.build_reports
    if hasattr(prb, "build_report"):
        return prb.build_report
    raise ImportError("ppt_report_builder.py must expose build_reports(...) (or build_report(...)).")


def _call_ppt_builder_robust(builder_fn, *, template_pptx_path: str, out_dir: str, analysis_id: str,
                            report: Dict[str, Any], raw: Dict[str, Any], export_pdf: bool):
    """
    Call the builder with multiple compatible signatures.

    We try:
      1) build_reports(template_pptx_path=..., out_dir=..., analysis_id=..., report_data={report,raw}, export_pdf=...)
      2) ... same but payload={report,raw}
      3) ... same but data={report,raw}
      4) ... passing report=..., raw=...
      5) positional fallback
    """
    payload = {"report": report, "raw": raw}
    sig = None
    try:
        sig = inspect.signature(builder_fn)
    except Exception:
        sig = None

    # Helper: build kwargs only if accepted
    def _filter_kwargs(kwargs: Dict[str, Any]) -> Dict[str, Any]:
        if sig is None:
            return kwargs
        params = sig.parameters
        return {k: v for k, v in kwargs.items() if k in params}

    # Attempt 1: report_data
    try:
        kwargs = _filter_kwargs({
            "template_pptx_path": template_pptx_path,
            "out_dir": out_dir,
            "analysis_id": analysis_id,
            "report_data": payload,
            "export_pdf": export_pdf,
        })
        return builder_fn(**kwargs)
    except TypeError:
        pass

    # Attempt 2: payload
    try:
        kwargs = _filter_kwargs({
            "template_pptx_path": template_pptx_path,
            "out_dir": out_dir,
            "analysis_id": analysis_id,
            "payload": payload,
            "export_pdf": export_pdf,
        })
        return builder_fn(**kwargs)
    except TypeError:
        pass

    # Attempt 3: data
    try:
        kwargs = _filter_kwargs({
            "template_pptx_path": template_pptx_path,
            "out_dir": out_dir,
            "analysis_id": analysis_id,
            "data": payload,
            "export_pdf": export_pdf,
        })
        return builder_fn(**kwargs)
    except TypeError:
        pass

    # Attempt 4: report + raw
    try:
        kwargs = _filter_kwargs({
            "template_pptx_path": template_pptx_path,
            "out_dir": out_dir,
            "analysis_id": analysis_id,
            "report": report,
            "raw": raw,
            "export_pdf": export_pdf,
        })
        return builder_fn(**kwargs)
    except TypeError:
        pass

    # Attempt 5: positional
    return builder_fn(template_pptx_path, out_dir, analysis_id, payload, export_pdf)


# -----------------------------
# PPT endpoint (JSON) INLINE
# Supports:
#   POST /v1/report/ppt              -> PPTX
#   POST /v1/report/ppt?format=pdf   -> PDF
# -----------------------------
@app.post("/v1/report/ppt")
async def report_ppt_v1(payload: Dict[str, Any] = Body(...), format: str = "pptx"):
    if not REPORTING_LIBS_AVAILABLE:
        return JSONResponse(status_code=503, content={"error": "PPT generation disabled (Missing dependencies)."})

    try:
        TEMPLATE_PPTX = _template_pptx_path()
        if not Path(TEMPLATE_PPTX).exists():
            return JSONResponse(
                status_code=500,
                content={"error": "PPT template not found", "template_path": TEMPLATE_PPTX, "cwd": os.getcwd()},
            )

        report = payload.get("report") or {}
        raw = payload.get("raw") or {}
        meta = report.get("meta") or {}
        analysis_id = str(meta.get("analysis_id") or report.get("analysis_id") or uuid.uuid4().hex[:10])

        out_dir = Path("outputs") / analysis_id
        out_dir.mkdir(parents=True, exist_ok=True)

        export_pdf = str(format).lower() == "pdf"

        builder_fn = _get_ppt_builder()
        pptx_path, pdf_path = _call_ppt_builder_robust(
            builder_fn,
            template_pptx_path=TEMPLATE_PPTX,
            out_dir=str(out_dir),
            analysis_id=analysis_id,
            report=report,
            raw=raw,
            export_pdf=export_pdf,
        )

        if export_pdf:
            if not pdf_path or not Path(pdf_path).exists():
                return JSONResponse(
                    status_code=500,
                    content={
                        "error": "PDF export failed or is not available in this environment.",
                        "hint": "Install LibreOffice (soffice) in the Render image to enable PPT→PDF conversion.",
                    },
                )
            return FileResponse(
                pdf_path,
                media_type="application/pdf",
                filename=f"ThermalAI_Report_{analysis_id}.pdf",
            )

        return FileResponse(
            pptx_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"ThermalAI_Report_{analysis_id}.pptx",
        )

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e), "traceback": traceback.format_exc()})


# -----------------------------
# Hotspot connected components + bounding boxes (PIL + numpy only)
# -----------------------------
def _connected_components_boxes(mask: np.ndarray, min_area_px: int = 200) -> List[Tuple[int, int, int, int]]:
    if mask is None:
        return []

    if mask.dtype != np.bool_:
        mask = mask.astype(bool)

    h, w = mask.shape
    visited = np.zeros((h, w), dtype=np.bool_)
    boxes: List[Tuple[int, int, int, int]] = []

    neighbors = [
        (-1, -1), (-1, 0), (-1, 1),
        (0, -1),           (0, 1),
        (1, -1),  (1, 0),  (1, 1),
    ]

    for y in range(h):
        for x in range(w):
            if not mask[y, x] or visited[y, x]:
                continue

            stack = [(y, x)]
            visited[y, x] = True

            min_x = max_x = x
            min_y = max_y = y
            area = 0

            while stack:
                cy, cx = stack.pop()
                area += 1

                min_x = min(min_x, cx)
                max_x = max(max_x, cx)
                min_y = min(min_y, cy)
                max_y = max(max_y, cy)

                for dy, dx in neighbors:
                    ny, nx = cy + dy, cx + dx
                    if 0 <= ny < h and 0 <= nx < w and mask[ny, nx] and not visited[ny, nx]:
                        visited[ny, nx] = True
                        stack.append((ny, nx))

            if area >= int(min_area_px):
                boxes.append((min_x, min_y, max_x, max_y))

    boxes.sort(key=lambda b: (b[2] - b[0] + 1) * (b[3] - b[1] + 1), reverse=True)
    return boxes


def draw_hotspot_boxes(
    pil_img: Image.Image,
    hotspot_mask: np.ndarray,
    min_area_px: int = 200,
    max_boxes: int = 20,
) -> Image.Image:
    img = pil_img.copy()
    draw = ImageDraw.Draw(img)

    boxes = _connected_components_boxes(hotspot_mask, min_area_px=min_area_px)[:max_boxes]
    outline = (255, 215, 0)
    thickness = 3

    for (x1, y1, x2, y2) in boxes:
        for t in range(thickness):
            draw.rectangle([x1 - t, y1 - t, x2 + t, y2 + t], outline=outline)

    return img



# ============================================================
# Billing / Entitlements (Stripe + Postgres)
# ============================================================
# This section enables:
# - Stripe Checkout (Project plans)
# - Stripe Webhooks (grant credits / activate subscription)
# - Server-side enforcement on /analyze and report downloads
#
# REQUIRED ENV VARS (Render -> Service -> Environment):
# - STRIPE_SECRET_KEY=sk_live_... (or sk_test_...)
# - STRIPE_WEBHOOK_SECRET=whsec_...
# - DATABASE_URL=postgres://... (Render Postgres internal URL)
#
# OPTIONAL:
# - VIP_EMAILS=comma,separated,emails
# - VIP_DOMAINS=comma,separated,domains (e.g. city.gov, utility.com)
#
# NOTE: This implementation uses `requests` (no stripe SDK) and `psycopg2` for Postgres.
# Add to requirements.txt on Render:
#   psycopg2-binary==2.9.*
#
# Stripe API docs: https://stripe.com/docs/api
# Webhook signature docs: https://stripe.com/docs/webhooks/signatures

BILLING_FREE_SCAN_LIMIT = 3  # Now per month (implemented in logic below)
BILLING_MONTHLY_QUOTA_DEFAULT = int(os.getenv("BILLING_MONTHLY_QUOTA_DEFAULT", "50"))

STRIPE_SECRET_KEY = os.getenv("STRIPE_SECRET_KEY", "")
STRIPE_WEBHOOK_SECRET = os.getenv("STRIPE_WEBHOOK_SECRET", "")
DATABASE_URL = os.getenv("DATABASE_URL", "")

VIP_EMAILS = {e.strip().lower() for e in os.getenv("VIP_EMAILS", "").split(",") if e.strip()}
VIP_EMAILS.add("jaime@allretech.org")  # Force add founder for dev/debug
VIP_DOMAINS = {d.strip().lower().lstrip("@") for d in os.getenv("VIP_DOMAINS", "").split(",") if d.strip()}

def _stripe_headers() -> Dict[str, str]:
    if not STRIPE_SECRET_KEY:
        raise RuntimeError("Missing STRIPE_SECRET_KEY env var.")
    return {"Authorization": f"Bearer {STRIPE_SECRET_KEY}"}

def _is_vip_email(email: Optional[str]) -> bool:
    print(f"DEBUG: _is_vip_email checking '{email}' against {VIP_EMAILS}")
    if not email:
        return False
    e = email.strip().lower()
    if e in VIP_EMAILS:
        return True
    if VIP_DOMAINS:
        domain = e.split("@")[-1] if "@" in e else ""
        if domain and domain.lower() in VIP_DOMAINS:
            return True
    return False

def _db_connect():
    """
    Returns (conn, kind) where kind is 'postgres' or 'sqlite'.
    Uses Postgres if DATABASE_URL is set and psycopg2 is available;
    otherwise falls back to local sqlite (best-effort).
    """
    if DATABASE_URL:
        try:
            import psycopg2  # type: ignore
            conn = psycopg2.connect(DATABASE_URL)
            return conn, "postgres"
        except ImportError:
            print("WARN: psycopg2 not found. Postgres unavailable, falling back to sqlite.")
        except Exception as e:
            print(f"WARN: Postgres connection failed: {e}, falling back to sqlite.")

    import sqlite3
    db_path = os.getenv("SQLITE_FALLBACK_PATH", "/tmp/thermalai_billing.sqlite")
    conn = sqlite3.connect(db_path)
    return conn, "sqlite"

def _db_exec(sql: str, params: Tuple = ()) -> None:
    conn, kind = _db_connect()
    try:
        cur = conn.cursor()
        # Normalize paramstyle: write SQL using '?' placeholders; convert for Postgres.
        if kind == "postgres":
            sql = sql.replace("?", "%s")
        cur.execute(sql, params)
        conn.commit()
    finally:
        try:
            conn.close()
        except Exception:
            pass

def _db_fetchone(sql: str, params: Tuple = ()) -> Optional[Tuple]:
    conn, kind = _db_connect()
    try:
        cur = conn.cursor()
        if kind == "postgres":
            sql = sql.replace("?", "%s")
        cur.execute(sql, params)
        row = cur.fetchone()
        return row
    finally:
        try:
            conn.close()
        except Exception:
            pass

def _db_fetchall(sql: str, params: Tuple = ()) -> List[Tuple]:
    conn, kind = _db_connect()
    try:
        cur = conn.cursor()
        if kind == "postgres":
            sql = sql.replace("?", "%s")
        cur.execute(sql, params)
        rows = cur.fetchall()
        return rows
    finally:
        try:
            conn.close()
        except Exception:
            pass

def _ensure_billing_tables() -> None:
    # Compatible schema for Postgres + SQLite
    _db_exec(
        """
        CREATE TABLE IF NOT EXISTS entitlements (
            user_id TEXT PRIMARY KEY,
            user_email TEXT,
            plan TEXT NOT NULL,
            downloads_allowed INTEGER NOT NULL DEFAULT 0,
            scan_credits_remaining INTEGER NOT NULL DEFAULT 0,
            monthly_quota INTEGER NOT NULL DEFAULT 0,
            monthly_used INTEGER NOT NULL DEFAULT 0,
            monthly_period_start TEXT,
            monthly_period_end TEXT,
            subscription_status TEXT,
            stripe_customer_id TEXT,
            stripe_subscription_id TEXT,
            free_scans_used INTEGER NOT NULL DEFAULT 0,
            vip_until TEXT,
            updated_at TEXT
        )
        """
    )

@app.on_event("startup")
def _billing_startup():
    try:
        _ensure_billing_tables()
        print("Billing tables ensured.")
    except Exception as e:
        print("WARN: Failed to ensure billing tables:", repr(e))

def _upsert_entitlement(
    user_id: str,
    user_email: Optional[str],
    plan: str,
    downloads_allowed: bool,
    scan_credits_remaining: int,
    monthly_quota: int,
    monthly_used: int,
    monthly_period_start: Optional[str],
    monthly_period_end: Optional[str],
    subscription_status: Optional[str],
    stripe_customer_id: Optional[str],
    stripe_subscription_id: Optional[str],
    free_scans_used: int,
    vip_until: Optional[str],
):
    now = datetime.utcnow().isoformat()
    # SQLite doesn't support TRUE/FALSE reliably -> store 0/1
    dl = 1 if downloads_allowed else 0

    # Upsert pattern compatible with Postgres + SQLite
    conn, kind = _db_connect()
    try:
        cur = conn.cursor()
        if kind == "postgres":
            cur.execute(
                """
                INSERT INTO entitlements (
                    user_id, user_email, plan, downloads_allowed,
                    scan_credits_remaining, monthly_quota, monthly_used,
                    monthly_period_start, monthly_period_end,
                    subscription_status, stripe_customer_id, stripe_subscription_id,
                    free_scans_used, vip_until, updated_at
                )
                VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
                ON CONFLICT (user_id) DO UPDATE SET
                    user_email=EXCLUDED.user_email,
                    plan=EXCLUDED.plan,
                    downloads_allowed=EXCLUDED.downloads_allowed,
                    scan_credits_remaining=EXCLUDED.scan_credits_remaining,
                    monthly_quota=EXCLUDED.monthly_quota,
                    monthly_used=EXCLUDED.monthly_used,
                    monthly_period_start=EXCLUDED.monthly_period_start,
                    monthly_period_end=EXCLUDED.monthly_period_end,
                    subscription_status=EXCLUDED.subscription_status,
                    stripe_customer_id=EXCLUDED.stripe_customer_id,
                    stripe_subscription_id=EXCLUDED.stripe_subscription_id,
                    free_scans_used=EXCLUDED.free_scans_used,
                    vip_until=EXCLUDED.vip_until,
                    updated_at=EXCLUDED.updated_at
                """,
                (
                    user_id, user_email, plan, dl,
                    scan_credits_remaining, monthly_quota, monthly_used,
                    monthly_period_start, monthly_period_end,
                    subscription_status, stripe_customer_id, stripe_subscription_id,
                    free_scans_used, vip_until, now
                ),
            )
        else:
            # SQLite: INSERT OR REPLACE
            cur.execute(
                """
                INSERT OR REPLACE INTO entitlements (
                    user_id, user_email, plan, downloads_allowed,
                    scan_credits_remaining, monthly_quota, monthly_used,
                    monthly_period_start, monthly_period_end,
                    subscription_status, stripe_customer_id, stripe_subscription_id,
                    free_scans_used, vip_until, updated_at
                )
                VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)
                """,
                (
                    user_id, user_email, plan, dl,
                    scan_credits_remaining, monthly_quota, monthly_used,
                    monthly_period_start, monthly_period_end,
                    subscription_status, stripe_customer_id, stripe_subscription_id,
                    free_scans_used, vip_until, now
                ),
            )
        conn.commit()
    finally:
        try:
            conn.close()
        except Exception:
            pass

def _get_entitlement(user_id: Optional[str], user_email: Optional[str]) -> Dict[str, Any]:
    # If no user context, treat as anonymous community (no persistence)
    if not user_id:
        return {
            "user_id": None,
            "plan": "community",
            "downloads_allowed": False,
            "scan_credits_remaining": 0,
            "monthly_quota": 0,
            "monthly_used": 0,
            "free_scans_used": 0,
            "allowed": True,
        }

    row = _db_fetchone(
        """
        SELECT user_id, user_email, plan, downloads_allowed, scan_credits_remaining,
               monthly_quota, monthly_used, monthly_period_start, monthly_period_end,
               subscription_status, stripe_customer_id, stripe_subscription_id,
               free_scans_used, vip_until
        FROM entitlements
        WHERE user_id = ?
        """,
        (user_id,),
    )

    # Auto-apply VIP by email (if configured)
    if _is_vip_email(user_email):
        # VIP: unlimited analysis + downloads (we implement as very large quota + credits)
        vip_until = None
        _upsert_entitlement(
            user_id=user_id,
            user_email=user_email,
            plan="vip",
            downloads_allowed=True,
            scan_credits_remaining=10_000_000,
            monthly_quota=10_000_000,
            monthly_used=0,
            monthly_period_start=None,
            monthly_period_end=None,
            subscription_status="vip",
            stripe_customer_id=None,
            stripe_subscription_id=None,
            free_scans_used=0,
            vip_until=vip_until,
        )
        row = _db_fetchone(
            """
            SELECT user_id, user_email, plan, downloads_allowed, scan_credits_remaining,
                   monthly_quota, monthly_used, monthly_period_start, monthly_period_end,
                   subscription_status, stripe_customer_id, stripe_subscription_id,
                   free_scans_used, vip_until
            FROM entitlements
            WHERE user_id = ?
            """,
            (user_id,),
        )

    if not row:
        # Create default community entitlement row
        _upsert_entitlement(
            user_id=user_id,
            user_email=user_email,
            plan="community",
            downloads_allowed=False,
            scan_credits_remaining=0,
            monthly_quota=0,
            monthly_used=0,
            monthly_period_start=None,
            monthly_period_end=None,
            subscription_status=None,
            stripe_customer_id=None,
            stripe_subscription_id=None,
            free_scans_used=0,
            vip_until=None,
        )
        row = _db_fetchone(
            """
            SELECT user_id, user_email, plan, downloads_allowed, scan_credits_remaining,
                   monthly_quota, monthly_used, monthly_period_start, monthly_period_end,
                   subscription_status, stripe_customer_id, stripe_subscription_id,
                   free_scans_used, vip_until
            FROM entitlements
            WHERE user_id = ?
            """,
            (user_id,),
        )

    (uid, uemail, plan, downloads_allowed, credits, m_quota, m_used,
     p_start, p_end, sub_status, cust_id, sub_id, free_used, vip_until) = row

    return {
        "user_id": uid,
        "user_email": uemail,
        "plan": plan,
        "downloads_allowed": bool(downloads_allowed),
        "scan_credits_remaining": int(credits or 0),
        "monthly_quota": int(m_quota or 0),
        "monthly_used": int(m_used or 0),
        "monthly_period_start": p_start,
        "monthly_period_end": p_end,
        "subscription_status": sub_status,
        "stripe_customer_id": cust_id,
        "stripe_subscription_id": sub_id,
        "free_scans_used": int(free_used or 0),
        "vip_until": vip_until,
    }

def _plan_allows_downloads(ent: Dict[str, Any]) -> bool:
    return bool(ent.get("downloads_allowed")) or ent.get("plan") in ("project", "enterprise", "vip")

def billing_can_analyze_internal(user_id: Optional[str], user_email: Optional[str], consume: bool) -> Dict[str, Any]:
    print(f"DEBUG: billing_can_analyze_internal called with user_id={user_id}, user_email={user_email}, consume={consume}")
    ent = _get_entitlement(user_id, user_email)

    # VIP always allowed
    if ent.get("plan") == "vip":
        ent["allowed"] = True
        ent["reason"] = None
        return ent

    plan = ent.get("plan", "community")

    # Community: 3 free scans PER MONTH
    if plan == "community":
        now_utc = datetime.now(timezone.utc)
        current_month_start = now_utc.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        
        # Check if we need to reset the monthly counter
        stored_start_iso = ent.get("monthly_period_start")
        monthly_used = int(ent.get("monthly_used", 0))

        # Parse stored date or default to None
        stored_start = None
        if stored_start_iso:
            try:
                stored_start = datetime.fromisoformat(stored_start_iso)
                # Handle offset-naive if necessary (though we use UTC usually)
                if stored_start.tzinfo is None:
                    stored_start = stored_start.replace(tzinfo=timezone.utc)
            except ValueError:
                pass

        # If stored start is different month than now, reset
        if not stored_start or (stored_start.year != current_month_start.year) or (stored_start.month != current_month_start.month):
            monthly_used = 0
            stored_start = current_month_start

        if monthly_used >= BILLING_FREE_SCAN_LIMIT:
            ent["allowed"] = False
            ent["reason"] = f"Monthly free limit reached ({BILLING_FREE_SCAN_LIMIT} analyses). Resets next month or upgrade to Project."
            return ent
        
        ent["allowed"] = True
        ent["reason"] = None
        
        if consume and user_id:
            monthly_used += 1
            # We also track lifetime usage just for stats, but logic depends on monthly_used
            free_used = int(ent.get("free_scans_used", 0)) + 1
            
            _upsert_entitlement(
                user_id=user_id,
                user_email=user_email,
                plan="community",
                downloads_allowed=False,
                scan_credits_remaining=0,
                monthly_quota=BILLING_FREE_SCAN_LIMIT,
                monthly_used=monthly_used,
                monthly_period_start=stored_start.isoformat(),
                monthly_period_end=None, # Not strictly needed for logic, but could compute end of month
                subscription_status=None,
                stripe_customer_id=ent.get("stripe_customer_id"),
                stripe_subscription_id=ent.get("stripe_subscription_id"),
                free_scans_used=free_used,
                vip_until=ent.get("vip_until"),
            )
            ent["monthly_used"] = monthly_used
        return ent

    # Project pack: consumes scan credits
    if plan == "project" and int(ent.get("scan_credits_remaining", 0)) > 0:
        ent["allowed"] = True
        ent["reason"] = None
        if consume and user_id:
            remaining = int(ent.get("scan_credits_remaining", 0)) - 1
            _upsert_entitlement(
                user_id=user_id,
                user_email=user_email,
                plan="project",
                downloads_allowed=True,
                scan_credits_remaining=max(0, remaining),
                monthly_quota=0,
                monthly_used=0,
                monthly_period_start=None,
                monthly_period_end=None,
                subscription_status=None,
                stripe_customer_id=ent.get("stripe_customer_id"),
                stripe_subscription_id=ent.get("stripe_subscription_id"),
                free_scans_used=int(ent.get("free_scans_used", 0)),
                vip_until=ent.get("vip_until"),
            )
            ent["scan_credits_remaining"] = max(0, remaining)
        return ent

    # Monthly: quota check within period
    if plan == "project_monthly":
        quota = int(ent.get("monthly_quota", BILLING_MONTHLY_QUOTA_DEFAULT))
        used = int(ent.get("monthly_used", 0))
        if used >= quota:
            ent["allowed"] = False
            ent["reason"] = f"Monthly quota reached ({quota} analyses)."
            return ent
        ent["allowed"] = True
        ent["reason"] = None
        if consume and user_id:
            used += 1
            _upsert_entitlement(
                user_id=user_id,
                user_email=user_email,
                plan="project_monthly",
                downloads_allowed=True,
                scan_credits_remaining=0,
                monthly_quota=quota,
                monthly_used=used,
                monthly_period_start=ent.get("monthly_period_start"),
                monthly_period_end=ent.get("monthly_period_end"),
                subscription_status=ent.get("subscription_status"),
                stripe_customer_id=ent.get("stripe_customer_id"),
                stripe_subscription_id=ent.get("stripe_subscription_id"),
                free_scans_used=int(ent.get("free_scans_used", 0)),
                vip_until=ent.get("vip_until"),
            )
            ent["monthly_used"] = used
        return ent

    # Project but no credits
    if plan == "project" and int(ent.get("scan_credits_remaining", 0)) <= 0:
        ent["allowed"] = False
        ent["reason"] = "No Project scan credits remaining. Buy more credits or upgrade."
        return ent

    # Enterprise: allow (you may enforce quota here later)
    if plan == "enterprise":
        ent["allowed"] = True
        ent["reason"] = None
        return ent

    # Default deny
    ent["allowed"] = False
    ent["reason"] = "Upgrade required."
    return ent

def _stripe_get_price_id_by_lookup_key(lookup_key: str) -> str:
    # GET /v1/prices?lookup_keys[]=...&active=true&limit=1
    url = "https://api.stripe.com/v1/prices"
    params = [("active", "true"), ("limit", "1"), ("lookup_keys[]", lookup_key)]
    r = requests.get(url, headers=_stripe_headers(), params=params, timeout=20)
    if r.status_code != 200:
        raise RuntimeError(f"Stripe prices lookup failed: {r.status_code} {r.text}")
    data = r.json()
    items = data.get("data", [])
    if not items:
        raise RuntimeError(f"No Stripe price found for lookup_key={lookup_key}")
    return items[0]["id"]

def _stripe_retrieve_subscription(subscription_id: str) -> Dict[str, Any]:
    url = f"https://api.stripe.com/v1/subscriptions/{subscription_id}"
    r = requests.get(url, headers=_stripe_headers(), timeout=20)
    if r.status_code != 200:
        raise RuntimeError(f"Stripe subscription retrieve failed: {r.status_code} {r.text}")
    return r.json()

def _grant_project_pack(user_id: str, user_email: Optional[str], credits_to_add: int, stripe_customer_id: Optional[str], source: str):
    ent = _get_entitlement(user_id, user_email)
    current = int(ent.get("scan_credits_remaining", 0))
    new_total = current + credits_to_add
    _upsert_entitlement(
        user_id=user_id,
        user_email=user_email,
        plan="project",
        downloads_allowed=True,
        scan_credits_remaining=new_total,
        monthly_quota=0,
        monthly_used=0,
        monthly_period_start=None,
        monthly_period_end=None,
        subscription_status=None,
        stripe_customer_id=stripe_customer_id or ent.get("stripe_customer_id"),
        stripe_subscription_id=None,
        free_scans_used=int(ent.get("free_scans_used", 0)),
        vip_until=ent.get("vip_until"),
    )

def _activate_monthly(user_id: str, user_email: Optional[str], subscription: Dict[str, Any], stripe_customer_id: Optional[str]):
    current_period_start = subscription.get("current_period_start")
    current_period_end = subscription.get("current_period_end")
    status = subscription.get("status")
    quota = BILLING_MONTHLY_QUOTA_DEFAULT

    _upsert_entitlement(
        user_id=user_id,
        user_email=user_email,
        plan="project_monthly",
        downloads_allowed=True,
        scan_credits_remaining=0,
        monthly_quota=quota,
        monthly_used=0,
        monthly_period_start=str(datetime.utcfromtimestamp(current_period_start).isoformat()) if current_period_start else None,
        monthly_period_end=str(datetime.utcfromtimestamp(current_period_end).isoformat()) if current_period_end else None,
        subscription_status=status,
        stripe_customer_id=stripe_customer_id,
        stripe_subscription_id=subscription.get("id"),
        free_scans_used=int(_get_entitlement(user_id, user_email).get("free_scans_used", 0)),
        vip_until=None,
    )

@app.get("/v1/billing/me")
async def billing_me(
    user_id: Optional[str] = None,
    user_email: Optional[str] = None,
    request: Request = None,
):
    # Allow headers as well (Base44 can send these)
    if request is not None:
        user_id = user_id or request.headers.get("X-User-Id")
        user_email = user_email or request.headers.get("X-User-Email")
    ent = _get_entitlement(user_id, user_email)
    ent["downloads_allowed"] = _plan_allows_downloads(ent)
    return ent

@app.post("/v1/billing/can_analyze")
async def billing_can_analyze(payload: Dict[str, Any] = Body(...)):
    user_id = payload.get("user_id")
    user_email = payload.get("user_email")
    consume = bool(payload.get("consume", False))
    ent = billing_can_analyze_internal(user_id=user_id, user_email=user_email, consume=consume)
    ent["downloads_allowed"] = _plan_allows_downloads(ent)
    return ent

@app.post("/v1/billing/checkout")
async def billing_checkout(payload: Dict[str, Any] = Body(...)):
    """
    Create a Stripe Checkout Session and return {url}.
    payload: { lookup_key, user_id, user_email, success_url, cancel_url }
    """
    lookup_key = payload.get("lookup_key")
    user_id = payload.get("user_id")
    user_email = payload.get("user_email")
    success_url = payload.get("success_url")
    cancel_url = payload.get("cancel_url")

    if not lookup_key:
        return JSONResponse(status_code=400, content={"error": "Missing lookup_key"})
    if not user_id or not user_email:
        return JSONResponse(status_code=400, content={"error": "Missing user_id or user_email"})
    if not success_url or not cancel_url:
        return JSONResponse(status_code=400, content={"error": "Missing success_url or cancel_url"})

    price_id = _stripe_get_price_id_by_lookup_key(lookup_key)
    mode = "subscription" if "monthly" in lookup_key else "payment"

    session_payload = {
        "mode": mode,
        "success_url": success_url,
        "cancel_url": cancel_url,
        "line_items[0][price]": price_id,
        "line_items[0][quantity]": 1,
        "client_reference_id": user_id,
        "customer_email": user_email,
        "metadata[user_id]": user_id,
        "metadata[user_email]": user_email,
        "metadata[lookup_key]": lookup_key,
    }

    # For subscriptions, copy metadata to subscription object as well
    if mode == "subscription":
        session_payload.update({
            "subscription_data[metadata][user_id]": user_id,
            "subscription_data[metadata][user_email]": user_email,
            "subscription_data[metadata][lookup_key]": lookup_key,
        })

    url = "https://api.stripe.com/v1/checkout/sessions"
    r = requests.post(url, headers=_stripe_headers(), data=session_payload, timeout=30)
    if r.status_code != 200:
        return JSONResponse(status_code=400, content={"error": f"Stripe error: {r.status_code}", "details": r.text})
    data = r.json()
    return {"url": data.get("url")}

def _verify_stripe_signature(payload: bytes, sig_header: str, secret: str, tolerance_sec: int = 300) -> bool:
    # Stripe-Signature: t=1492774577,v1=...,v0=...
    try:
        parts = {kv.split("=", 1)[0]: kv.split("=", 1)[1] for kv in sig_header.split(",")}
        timestamp = int(parts.get("t", "0"))
        signature = parts.get("v1", "")
        if not signature or not timestamp:
            return False

        if abs(int(time.time()) - timestamp) > tolerance_sec:
            return False

        signed_payload = f"{timestamp}.{payload.decode('utf-8')}".encode("utf-8")
        expected = hmac.new(secret.encode("utf-8"), signed_payload, hashlib.sha256).hexdigest()
        # Constant-time compare
        return hmac.compare_digest(expected, signature)
    except Exception:
        return False

@app.post("/v1/billing/webhook")
async def billing_webhook(request: Request):
    if not STRIPE_WEBHOOK_SECRET:
        return JSONResponse(status_code=500, content={"error": "Missing STRIPE_WEBHOOK_SECRET env var"})

    raw = await request.body()
    sig = request.headers.get("Stripe-Signature", "")

    if not _verify_stripe_signature(raw, sig, STRIPE_WEBHOOK_SECRET):
        return JSONResponse(status_code=400, content={"error": "Invalid Stripe signature"})

    event = json.loads(raw.decode("utf-8"))
    event_type = event.get("type")
    obj = (event.get("data") or {}).get("object") or {}

    try:
        if event_type == "checkout.session.completed":
            # Fulfillment: grant credits / activate subscription
            lookup_key = (obj.get("metadata") or {}).get("lookup_key")
            user_id = (obj.get("metadata") or {}).get("user_id") or obj.get("client_reference_id")
            user_email = (obj.get("metadata") or {}).get("user_email") or obj.get("customer_details", {}).get("email") or obj.get("customer_email")
            stripe_customer_id = obj.get("customer")
            subscription_id = obj.get("subscription")

            if not (lookup_key and user_id):
                return {"status": "ignored"}

            if "monthly" in lookup_key and subscription_id:
                sub = _stripe_retrieve_subscription(subscription_id)
                _activate_monthly(user_id=user_id, user_email=user_email, subscription=sub, stripe_customer_id=stripe_customer_id)
            else:
                # One-time packs
                credits_map = {
                    "project_scan_1": 1,
                    "project_pack_10": 10,
                    "project_pack_50": 50,
                }
                credits = credits_map.get(lookup_key, 0)
                if credits > 0:
                    _grant_project_pack(user_id=user_id, user_email=user_email, credits_to_add=credits, stripe_customer_id=stripe_customer_id, source=lookup_key)

        elif event_type in ("customer.subscription.updated", "customer.subscription.deleted"):
            sub = obj
            user_id = (sub.get("metadata") or {}).get("user_id")
            user_email = (sub.get("metadata") or {}).get("user_email")
            status = sub.get("status")
            stripe_customer_id = sub.get("customer")
            if user_id:
                if status in ("active", "trialing"):
                    _activate_monthly(user_id=user_id, user_email=user_email, subscription=sub, stripe_customer_id=stripe_customer_id)
                else:
                    # Downgrade: keep plan but disable usage (or set back to community)
                    _upsert_entitlement(
                        user_id=user_id,
                        user_email=user_email,
                        plan="community",
                        downloads_allowed=False,
                        scan_credits_remaining=0,
                        monthly_quota=0,
                        monthly_used=0,
                        monthly_period_start=None,
                        monthly_period_end=None,
                        subscription_status=status,
                        stripe_customer_id=stripe_customer_id,
                        stripe_subscription_id=sub.get("id"),
                        free_scans_used=int(_get_entitlement(user_id, user_email).get("free_scans_used", 0)),
                        vip_until=None,
                    )

        elif event_type == "invoice.payment_failed":
            # Optional: mark subscription as past_due -> lock monthly usage
            inv = obj
            subscription_id = inv.get("subscription")
            if subscription_id:
                sub = _stripe_retrieve_subscription(subscription_id)
                user_id = (sub.get("metadata") or {}).get("user_id")
                user_email = (sub.get("metadata") or {}).get("user_email")
                if user_id:
                    _upsert_entitlement(
                        user_id=user_id,
                        user_email=user_email,
                        plan="community",
                        downloads_allowed=False,
                        scan_credits_remaining=0,
                        monthly_quota=0,
                        monthly_used=0,
                        monthly_period_start=None,
                        monthly_period_end=None,
                        subscription_status="past_due",
                        stripe_customer_id=sub.get("customer"),
                        stripe_subscription_id=sub.get("id"),
                        free_scans_used=int(_get_entitlement(user_id, user_email).get("free_scans_used", 0)),
                        vip_until=None,
                    )

        # invoice.payment_succeeded can be used for analytics; monthly activation is handled above.

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": "Webhook handler failed", "details": repr(e)})

    return {"status": "ok"}


# -----------------------------
# Analyze
# -----------------------------
@app.post("/analyze")
async def analyze(
    request: Request,

    # billing / entitlement context (passed by frontend)
    user_id: Optional[str] = Form(default=None),
    user_email: Optional[str] = Form(default=None),

    thermal_image: Optional[UploadFile] = File(default=None),
    rgb_image: Optional[UploadFile] = File(default=None),

    city: Optional[str] = Form(default=None),
    country: Optional[str] = Form(default=None),
    latitude: Optional[str] = Form(default=None),
    longitude: Optional[str] = Form(default=None),
    datetime_iso: Optional[str] = Form(default=None),

    t_inside: str = Form(default="22"),
    t_outside: Optional[str] = Form(default=None),

    heating_base_temp_c: str = Form(default="13"),

    fuel_price_eur_per_kwh: str = Form(default="0.12"),
    discount_rate: Optional[str] = Form(default=None),
    inflation_rate: Optional[str] = Form(default=None),

    facade_area_m2: Optional[str] = Form(default=None),
    hotspot_area_m2_override: Optional[str] = Form(default=None),

    # building characteristics
    building_type: Optional[str] = Form(default=None),
    building_year: Optional[str] = Form(default=None),
    floor_area_m2: Optional[str] = Form(default=None),
    envelope_area_m2: Optional[str] = Form(default=None),
    num_stories: Optional[str] = Form(default=None),
    heating_system: Optional[str] = Form(default=None),
    climate_zone: Optional[str] = Form(default=None),
    hdd: Optional[str] = Form(default=None),

    # environmental conditions
    outdoor_rh_percent: Optional[str] = Form(default=None),
    wind_speed_mps: Optional[str] = Form(default=None),
    sky_conditions: Optional[str] = Form(default=None),

    # address/meta
    address: Optional[str] = Form(default=None),
    google_maps_link: Optional[str] = Form(default=None),
    important_note: Optional[str] = Form(default=None),

    # U-values
    u_current_wall: Optional[str] = Form(default=None),
    u_improved_wall: Optional[str] = Form(default=None),
    u_current_window: Optional[str] = Form(default=None),
    u_improved_window: Optional[str] = Form(default=None),
    u_current_door: Optional[str] = Form(default=None),
    u_improved_door: Optional[str] = Form(default=None),

    # materials
    material_current_wall: Optional[str] = Form(default="uninsulated_brick_wall"),
    material_improved_wall: Optional[str] = Form(default="insulated_wall"),
    material_current_window: Optional[str] = Form(default="single_glazed_window"),
    material_improved_window: Optional[str] = Form(default="double_glazed_window"),
    material_current_door: Optional[str] = Form(default="default"),
    material_improved_door: Optional[str] = Form(default="default"),

    overlay_threshold_percentile: str = Form(default="95"),

    include_gamma_payload: str = Form(default="false"),
    include_overlay_base64: str = Form(default="true"),
    auto_register: str = Form(default="true"),  # NEW: allow disabling registration if needed
):
    # Check if ML/Image libs are available
    if not (ML_AVAILABLE or CV2_AVAILABLE) or not IMAGE_LIBS_AVAILABLE:
        # Return a mock response if ML libs are missing (Vercel LITE mode)
        print("Warning: ML/Image libs missing but proceeding to try analysis...")
        pass
        # msg = "AI analysis unavailable."
        # if not IMAGE_LIBS_AVAILABLE:
        #     msg += " (Image libs missing)"
            
        # print(f"Analyze called but libraries missing. Returning mock. ML={ML_AVAILABLE}, IMG={IMAGE_LIBS_AVAILABLE}")
        # return JSONResponse(content={
        #     "status": "success",
        #     "mock": True,
        #     "message": msg,
        #     "analysis_result": {
        #         "summary": "AI processing is disabled in this environment.",
        #         "issues": []
        #     },
        #     "artifacts": {}
        # })

    if SEG_MODEL is None:
        return JSONResponse(status_code=503, content={"error": "Model not loaded yet. Please retry in a few seconds."})

    # -----------------------------
    # Entitlement enforcement (server-side)
    # -----------------------------
    ent = billing_can_analyze_internal(user_id=user_id, user_email=user_email, consume=True)
    if not ent["allowed"]:
        return JSONResponse(
            status_code=402,
            content={"error": ent.get("reason", "Upgrade required"), "entitlement": ent},
        )

    # Proxy fallback: try request.form()
    if rgb_image is None or thermal_image is None:
        form = await request.form()
        rgb_image = rgb_image or form.get("rgb_image")
        thermal_image = thermal_image or form.get("thermal_image")
        if rgb_image is None or thermal_image is None:
            return JSONResponse(
                status_code=400,
                content={
                    "error": "Missing files. Expected multipart form-data fields rgb_image and thermal_image.",
                    "received_keys": list(form.keys()),
                },
            )

    include_gamma = _safe_bool(include_gamma_payload, default=False)
    include_overlay = _safe_bool(include_overlay_base64, default=True)

    # Safe defaults
    T_inside = float(t_inside) if t_inside else 22.0
    fuel_price = float(fuel_price_eur_per_kwh) if fuel_price_eur_per_kwh else 0.12
    base_temp = float(heating_base_temp_c) if heating_base_temp_c else 13.0
    overlay_pct = float(overlay_threshold_percentile) if overlay_threshold_percentile else 95.0

    # Always default to 3% if not provided
    dr = _safe_float(discount_rate)
    ir = _safe_float(inflation_rate)
    dr = 0.03 if dr is None else float(dr)
    ir = 0.03 if ir is None else float(ir)

    lat = _safe_float(latitude)
    lon = _safe_float(longitude)

    rgb_bytes = await rgb_image.read()
    thermal_bytes = await thermal_image.read()

    vis_img = Image.open(io.BytesIO(rgb_bytes)).convert("RGB")
    thr_img = Image.open(io.BytesIO(thermal_bytes)).convert("RGB")

    vis_img = _resize_max(vis_img, max_side=1024)
    thr_img = _resize_max(thr_img, max_side=1024)

    # >>> thermal -> RGB registration (audited engine) with interpolation support mask
    if _safe_bool(auto_register, default=True) and CV2_AVAILABLE:
        thr_img, valid_support_mask, registration_meta = register_thermal_to_rgb(vis_img, thr_img)
    else:
        # registration disabled by caller: plain resize, whole frame treated as supported
        thr_img = thr_img.resize(vis_img.size, Image.Resampling.BILINEAR)
        valid_support_mask = np.ones((vis_img.size[1], vis_img.size[0]), dtype=bool)
        registration_meta = {
            "used": False, "method": "resize_only", "confidence": None, "reliable": False,
            "quality_label": "poor", "support_coverage": 1.0, "rejections": ["auto_register disabled"],
        }

    # Segmentation; the active analysis region is facade AND valid interpolation support.
    # Threshold statistics are computed strictly inside this region (never on the whole frame).
    seg = SEG_MODEL.predict_masks(vis_img)
    facade_mask = (seg.wall_mask | seg.window_mask | seg.door_mask) & valid_support_mask
    hs = detect_hotspot_mask(thr_img, threshold_percentile=overlay_pct, active_analysis_mask=facade_mask)
    analysis_flags: List[str] = []
    if hs.status == "empty_active_mask":
        analysis_flags.append("empty_active_region: no facade pixels with valid thermal support; hotspot statistics not computed")
    elif hs.status == "flat_image":
        analysis_flags.append("flat_thermal_image: no thermal contrast inside the active region")
    if not registration_meta.get("reliable", False):
        analysis_flags.append(f"registration_not_verified: method={registration_meta.get('method')} ({registration_meta.get('quality_label')})")

    # Artifacts
    overlay_b64 = None
    rgb_b64 = None
    rgb_boxes_b64 = None
    thermal_boxes_b64 = None

    if include_overlay:
        overlay_img = overlay_mask_on_rgb(vis_img, hs.mask)
        overlay_b64 = encode_image_to_base64_png(overlay_img)
        rgb_b64 = encode_image_to_base64_png(vis_img)

        rgb_boxes_img = draw_hotspot_boxes(vis_img, hs.mask, min_area_px=200, max_boxes=20)
        thr_boxes_img = draw_hotspot_boxes(thr_img, hs.mask, min_area_px=200, max_boxes=20)
        rgb_boxes_b64 = encode_image_to_base64_png(rgb_boxes_img)
        thermal_boxes_b64 = encode_image_to_base64_png(thr_boxes_img)
        thermal_b64 = encode_image_to_base64_png(thr_img)

    # Outdoor temp: user > API > fallback
    T_outside = _safe_float(t_outside)
    t_out_src = "user_input"
    if T_outside is None:
        fetched = get_outdoor_temperature_c(
            city=city, country=country, latitude=lat, longitude=lon, datetime_iso=datetime_iso
        )
        if fetched is not None:
            T_outside = float(fetched)
            t_out_src = "weather_api"
        else:
            T_outside = 5.0
            t_out_src = "fallback_default_5C"

    delta_t_capture = max(0.0, T_inside - T_outside)

    deg_hours = degree_hours_below_base(
        base_temp, city=city, country=country, latitude=lat, longitude=lon, prefer_city_table=True
    )
    deg_src = "city_table_or_openmeteo"
    if deg_hours is None:
        deg_hours = 30000.0
        deg_src = "fallback_default_30000"

    facade_area = _safe_float(facade_area_m2)
    hotspot_override = _safe_float(hotspot_area_m2_override)

    total_pixels = seg.counts["wall_pixels"] + seg.counts["window_pixels"] + seg.counts["door_pixels"]
    total_pixels = max(total_pixels, 1)

    if facade_area is not None:
        comp_area = {
            "wall": facade_area * (seg.counts["wall_pixels"] / total_pixels),
            "window": facade_area * (seg.counts["window_pixels"] / total_pixels),
            "door": facade_area * (seg.counts["door_pixels"] / total_pixels),
        }
        comp_area_src = "scaled_from_facade_area_and_segmentation"
    else:
        comp_area = {
            "wall": seg.counts["wall_pixels"] / total_pixels,
            "window": seg.counts["window_pixels"] / total_pixels,
            "door": seg.counts["door_pixels"] / total_pixels,
        }
        comp_area_src = "relative_only_no_facade_area"

    # --- DEBUG: HOTSPOT STATS ---
    try:
        hs_pixels = int(hs.mask.sum())
        sys.stderr.write(f"DEBUG: Hotspot detection: {hs_pixels} pixels (threshold p{overlay_pct})\n")
        
        # Check raw boxes before drawing
        sys.stderr.write("DEBUG: Generating boxes...\n")
        raw_boxes = _connected_components_boxes(hs.mask, min_area_px=200)
        sys.stderr.write(f"DEBUG: Found {len(raw_boxes)} boxes.\n")
    except Exception as e:
        sys.stderr.write(f"DEBUG: Failed to log hotspot stats: {e}\n")
    # ----------------------------

    # External combined heat-transfer coefficient from the reported wind speed and the
    # outdoor air temperature (surface assumed at outdoor air temperature for screening).
    wind_for_h = _safe_float(wind_speed_mps)
    wind_for_h = DEFAULT_WIND_SPEED_M_S if wind_for_h is None else max(0.1, wind_for_h)
    h_ext = calculate_external_heat_transfer_coefficient(
        wind_speed_m_s=wind_for_h, surface_temp_k=float(T_outside) + 273.15
    )
    if delta_t_capture < MIN_RELIABLE_DELTA_T_C:
        analysis_flags.append(
            f"low_delta_t_capture: indoor-outdoor difference {delta_t_capture:.1f} K below {MIN_RELIABLE_DELTA_T_C:.0f} K; annualised proxy unreliable"
        )

    # Component masks restricted to the same active region used for the statistics
    masks = {
        "wall": seg.wall_mask & valid_support_mask,
        "window": seg.window_mask & valid_support_mask,
        "door": seg.door_mask & valid_support_mask,
    }

    components: Dict[str, Any] = {}
    totals = {
        "instantaneous_watts": 0.0,
        "annual_kwh_delta": 0.0,
        "annual_cost_delta": 0.0,
        "annual_kwh_u": 0.0,
        "annual_cost_u": 0.0,
        "annual_kwh_theoretical": 0.0,
        "annual_cost_theoretical": 0.0,
    }

    for name, cmask in masks.items():
        comp_pixels = int(cmask.sum())
        if comp_pixels <= 0:
            hs_in_comp = 0
            hs_ratio_in_comp = 0.0
        else:
            hs_in_comp = int((hs.mask & cmask).sum())
            hs_ratio_in_comp = float(hs_in_comp) / float(comp_pixels)

        if hotspot_override is not None and name == "wall":
            hs_area_m2 = float(hotspot_override)
            hs_area_src = "user_override_for_wall"
        else:
            hs_area_m2 = float(comp_area[name]) * float(hs_ratio_in_comp)
            hs_area_src = "component_area_x_hotspot_ratio"

        inst_w = instantaneous_loss_proxy_watts(
            hs_area_m2, T_inside, T_outside, kappa_calibration=DEFAULT_KAPPA_CALIBRATION, h_ext=h_ext
        )
        ann = annualize_proxy_kwh(inst_w, delta_t_capture, deg_hours)
        annual_kwh_delta = ann.annual_kwh
        if not ann.valid and ann.status != "zero_loss":
            for w in (ann.warnings or [ann.status]):
                msg = f"{name}: {w}"
                if msg not in analysis_flags:
                    analysis_flags.append(msg)

        if name == "wall":
            u_cur = _safe_float(u_current_wall) or infer_u_value(material_current_wall)
            u_imp = _safe_float(u_improved_wall) or infer_u_value(material_improved_wall)
        elif name == "window":
            u_cur = _safe_float(u_current_window) or infer_u_value(material_current_window)
            u_imp = _safe_float(u_improved_window) or infer_u_value(material_improved_window)
        else:
            u_cur = _safe_float(u_current_door) or infer_u_value(material_current_door)
            u_imp = _safe_float(u_improved_door) or infer_u_value(material_improved_door)

        hdd_equiv = float(deg_hours) / 24.0
        # Savings (Retrofit Potential) - based on Area
        # We assume retrofit applies to the WHOLE component area if we want to be realistic about "savings potential"
        # BUT originally this was likely just hotspot area.
        # Let's align savings with the new "Total Theoretical" approach -> Savings on the whole area.
        comp_total_area = float(comp_area[name])
        annual_kwh_u = annual_kwh_saved_u_method(u_cur, u_imp, comp_total_area, hdd_equiv)

        # Total Theoretical Loss (Current State)
        annual_kwh_theoretical = annual_total_loss_u_method(u_cur, comp_total_area, hdd_equiv)

        annual_cost_delta = round(annual_kwh_delta * fuel_price, 2)
        annual_cost_u = round(annual_kwh_u * fuel_price, 2)
        annual_cost_theoretical = round(annual_kwh_theoretical * fuel_price, 2)

        components[name] = {
            "component_pixels": comp_pixels,
            "hotspot_pixels_in_component": hs_in_comp,
            "hotspot_ratio_in_component": round(hs_ratio_in_comp, 6),
            "hotspot_area_m2": round(hs_area_m2, 4),
            "component_total_area_m2": round(comp_total_area, 4),
            "hotspot_area_source": hs_area_src,
            "instantaneous_watts": round(inst_w, 4),
            "annual_kwh_hotspot_delta": round(annual_kwh_delta, 4),
            "annual_kwh_savings_potential": round(annual_kwh_u, 4),
            "annual_kwh_total_loss": round(annual_kwh_theoretical, 4)
        }

        totals["instantaneous_watts"] += inst_w
        totals["annual_kwh_delta"] += annual_kwh_delta
        totals["annual_cost_delta"] += annual_cost_delta
        totals["annual_kwh_u"] += annual_kwh_u
        totals["annual_cost_u"] += annual_cost_u
        totals["annual_kwh_theoretical"] += annual_kwh_theoretical
        totals["annual_cost_theoretical"] += annual_cost_theoretical

    totals = {k: (round(v, 4) if "kwh" in k or "watts" in k else round(v, 2)) for k, v in totals.items()}
    # Use the Theoretical Total Cost for the multi-year projection
    # Multi-year projection: contract is (annual kWh, price per kWh, inflation, discount)
    totals["multi_year_costs_delta"] = compute_multi_year_costs(
        totals["annual_kwh_theoretical"], energy_price_per_kwh=fuel_price, inflation_rate=ir, discount_rate=dr
    )
    totals["multi_year_costs_hotspot_proxy"] = compute_multi_year_costs(
        totals["annual_kwh_delta"], energy_price_per_kwh=fuel_price, inflation_rate=ir, discount_rate=dr
    )
    totals["physics_assumptions"] = {
        "kappa_calibration": DEFAULT_KAPPA_CALIBRATION,
        "h_ext_w_per_m2k": round(h_ext, 3),
        "wind_speed_m_s_used": wind_for_h,
        "delta_t_capture_c": round(delta_t_capture, 3),
        "hotspot_status": hs.status,
        "hotspot_threshold_gray": round(float(hs.threshold), 2),
        "active_region_pixels": int(hs.total_pixels),
    }

    analysis_id = uuid.uuid4().hex[:10]
    api_base = "" # Let frontend handle base URL via proxy

    response: Dict[str, Any] = {
        "analysis_id": analysis_id,
        "inputs": {
            "api_base": api_base,
            "city": city,
            "country": country,
            "latitude": lat,
            "longitude": lon,
            "datetime_iso": datetime_iso,
            "t_inside_c": T_inside,
            "t_outside_c": T_outside,
            "t_outside_source": t_out_src,
            "heating_base_temp_c": base_temp,
            "degree_hours_annual": round(float(deg_hours), 2),
            "degree_hours_source": deg_src,
            "fuel_price_eur_per_kwh": fuel_price,
            "discount_rate": dr,
            "inflation_rate": ir,
            "overlay_threshold_percentile": overlay_pct,
            "facade_area_m2": facade_area,
            "component_area_source": comp_area_src,
            "segmentation_counts": seg.counts,

            "building_type": building_type,
            "building_year": building_year,
            "floor_area_m2": _safe_float(floor_area_m2),
            "envelope_area_m2": _safe_float(envelope_area_m2),
            "num_stories": _safe_float(num_stories),
            "heating_system": heating_system,
            "climate_zone": climate_zone,
            "hdd": _safe_float(hdd),

            "outdoor_rh_percent": _safe_float(outdoor_rh_percent),
            "wind_speed_mps": _safe_float(wind_speed_mps),
            "sky_conditions": sky_conditions,

            "address": address,
            "google_maps_link": google_maps_link,
            "important_note": important_note,

            # registration diagnostics + analysis quality flags visible to frontend / report
            "registration": registration_meta,
            "analysis_flags": analysis_flags,
            "backend_version": BACKEND_VERSION,
            "registration_engine_version": REGISTRATION_ENGINE_VERSION,
        },
        "results": {"components": components, "totals": totals},
    }

    if include_overlay:
        response["artifacts"] = {
            "overlay_image_base64_png": overlay_b64,
            "rgb_image_base64_png": rgb_b64,
            
            # Frontend keys (Results.jsx expects specific names)
            "boxed_rgb_image_base64_png": rgb_boxes_b64,
            "rgb_boxed_image_base64_png": rgb_boxes_b64, # match Results.jsx
            "thermal_image_base64_png": thermal_b64,
            "thermal_boxed_image_base64_png": thermal_boxes_b64, # match Results.jsx

            # PPT Builder keys (legacy/backend support)
            "rgb_hotspot_boxes_base64_png": rgb_boxes_b64,
            "thermal_hotspot_boxes_base64_png": thermal_boxes_b64,
        }

    if include_gamma:
        response["gamma_payload"] = build_gamma_payload(response)

    report = build_report(response)

    report_meta = report.get("meta") or {}
    report_meta.update({
        "analysis_id": analysis_id,
        "api_base": api_base,

        "building_type": building_type,
        "building_year": building_year,
        "floor_area_m2": _safe_float(floor_area_m2),
        "envelope_area_m2": _safe_float(envelope_area_m2),
        "num_stories": _safe_float(num_stories),
        "heating_system": heating_system,
        "climate_zone": climate_zone,
        "hdd": _safe_float(hdd),

        "outdoor_rh_percent": _safe_float(outdoor_rh_percent),
        "wind_speed_mps": _safe_float(wind_speed_mps),
        "sky_conditions": sky_conditions,

        "address": address,
        "google_maps_link": google_maps_link,
        "important_note": important_note,

        "discount_rate": dr,
        "inflation_rate": ir,
        "overlay_threshold_percentile": overlay_pct,
        "date_of_report": datetime.now().strftime("%d/%m/%Y"),

        # propagate registration meta into meta as well if you want to use it in PPT
        "registration": registration_meta,
        "analysis_flags": analysis_flags,
        "version": BACKEND_VERSION,
    })
    report["meta"] = report_meta

    return JSONResponse(content={"report": report, "raw": response})


# -----------------------------
# Dashboard / Persistence
# -----------------------------
DATA_DIR = Path("data")
DATA_DIR.mkdir(exist_ok=True)
DASHBOARD_FILE = DATA_DIR / "user_dashboards.json"

class DashboardStore:
    def __init__(self):
        self.lock = threading.Lock()
        self._load()

    def _load(self):
        if DASHBOARD_FILE.exists():
            try:
                with open(DASHBOARD_FILE, "r") as f:
                    self.data = json.load(f)
            except Exception:
                print("Warning: Corrupt dashboard file, starting fresh.")
                self.data = {}
        else:
            self.data = {}
            # Seed Gyor Pilot for Jaime
            self._seed_pilot()
            
    def _seed_pilot(self):
        # We can seed the initial student hostel for everyone or just Jaime
        # For now, let's leave it empty and let frontend combine hardcoded + dynamic
        # Or better: Frontend hardcodes "Public" pilots, Backend stores "User" additions.
        pass

    def _save(self):
        with open(DASHBOARD_FILE, "w") as f:
            json.dump(self.data, f, indent=2)

    def get_user_buildings(self, email: str) -> List[Dict]:
        email = email.lower().strip()
        # Reload to ensure we see data from other workers/processes (if any)
        self._load()
        with self.lock:
            data = self.data.get(email, [])
            print(f"DEBUG: get_user_buildings for '{email}' -> found {len(data)}")
            return data

    def add_building(self, email: str, building: Dict):
        email = email.lower().strip()
        with self.lock:
            # Ensure we have latest before appending
            self._load()
            
            if email not in self.data:
                self.data[email] = []
            
            # Simple dedup by ID
            if "id" not in building:
                building["id"] = str(uuid.uuid4())
            
            # Check if exists
            exists = any(b.get("id") == building["id"] for b in self.data[email])
            if not exists:
                self.data[email].append(building)
                print(f"DEBUG: Added building {building['id']} for '{email}'")
                self._save()
            else:
                print(f"DEBUG: Building {building['id']} already exists for '{email}'")
            
            return building

dashboard_store = DashboardStore()

@app.get("/v1/dashboard")
async def get_dashboard(x_user_email: Optional[str] = Header(None)):
    if not x_user_email:
        # Return empty or error? For now empty list if no auth
        return []
    return dashboard_store.get_user_buildings(x_user_email)

class AddBuildingRequest(BaseModel):
    user_email: EmailStr
    building: Dict[str, Any]

@app.post("/v1/dashboard/add")
async def add_to_dashboard(req: AddBuildingRequest):
    # Validate location
    b = req.building
    has_loc = (
        (b.get('lat') and b.get('lng')) or 
        b.get('google_maps_link') or 
        (b.get('address') and len(b.get('address', '')) > 5)
    )
    
    if not has_loc:
        raise HTTPException(status_code=400, detail="Cannot add to dashboard: Missing Latitude/Longitude, Google Maps Link, or Full Address.")

    dashboard_store.add_building(req.user_email, b)
    return {"status": "ok", "building": b}

class DeleteBuildingRequest(BaseModel):
    user_email: EmailStr
    building_id: str

@app.post("/v1/dashboard/delete")
async def delete_from_dashboard(req: DeleteBuildingRequest):
    # Implement delete logic in store
    with dashboard_store.lock:
        if req.user_email in dashboard_store.data:
            initial_len = len(dashboard_store.data[req.user_email])
            dashboard_store.data[req.user_email] = [
                b for b in dashboard_store.data[req.user_email] 
                if b.get("id") != req.building_id
            ]
            if len(dashboard_store.data[req.user_email]) < initial_len:
                dashboard_store._save()
                return {"status": "ok", "deleted_id": req.building_id}
    
    raise HTTPException(status_code=404, detail="Building not found or not authorized.")

@app.post("/v1/dashboard/wipe_user")
async def wipe_user_dashboard(user_email: str = Body(..., embed=True)):
    with dashboard_store.lock:
        if user_email in dashboard_store.data:
            dashboard_store.data[user_email] = []
            dashboard_store._save()
            return {"status": "ok", "message": f"Wiped data for {user_email}"}
    return {"status": "ok", "message": "User not found, nothing to wipe"}

# ============================================================
# SPA Catch-All (Moved to end to prevent shadowing API routes)
# ============================================================
if frontend_dist.exists():
    @app.get("/{full_path:path}")
    async def serve_spa(full_path: str):
        # Allow API calls to pass through (already handled by routers above, but safe check)
        if full_path.startswith("api/") or full_path.startswith("docs") or full_path.startswith("openapi") or full_path.startswith("v1/"):
            raise HTTPException(status_code=404, detail="Not Found")
            
        # Check if a static file exists at this path in frontend/dist
        potential_path = frontend_dist / full_path
        if potential_path.exists() and potential_path.is_file():
            return FileResponse(potential_path)

        # Serve index.html for any other route (SPA Fallback)
        index_path = frontend_dist / "index.html"
        return FileResponse(index_path)
else:
    print(f"WARNING: Frontend dist not found at {frontend_dist}. API only mode.")
    @app.get("/")
    def root() -> Dict[str, Any]:
        return {"status": "ok", "service": "ThermalAI API (Frontend not attached)"}
