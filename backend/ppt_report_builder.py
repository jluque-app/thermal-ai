# ppt_report_builder.py
from __future__ import annotations

import base64
import math
import re
import shutil
import subprocess
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Optional, Tuple, List

from pptx import Presentation
from pptx.util import Pt

# -----------------------------
# Token parsing + defaults
# -----------------------------
# Support {{TOKEN}}, {{ TOKEN }}, and tokens split across runs.
_TOKEN_PATTERN = re.compile(r"\{\{\s*([A-Za-z0-9_]+)\s*\}\}")
_MISSING_TEXT = "n.a."


# -----------------------------
# Small helpers
# -----------------------------
def _dig(d: Any, path: str, default=None):
    """Safe nested getter: _dig(obj, 'a.b.c') or _dig(obj, 'a[0].b')."""
    try:
        cur = d
        for part in path.replace("]", "").split("."):
            if not part:
                continue
            if "[" in part:
                key, idx = part.split("[", 1)
                if key:
                    cur = cur.get(key, {})
                cur = cur[int(idx)]
            else:
                if isinstance(cur, dict):
                    cur = cur.get(part, default)
                else:
                    return default
        return cur if cur is not None else default
    except Exception:
        return default


def _safe_float(x: Any) -> Optional[float]:
    try:
        if x is None:
            return None
        if isinstance(x, (int, float)):
            v = float(x)
            if math.isnan(v) or math.isinf(v):
                return None
            return v
        s = str(x).strip()
        if not s:
            return None
        s = s.replace("€", "").replace("$", "").replace("\u202f", "").replace(" ", "")
        # allow either comma or dot decimal
        if s.count(",") == 1 and s.count(".") == 0:
            s = s.replace(",", ".")
        # remove thousands separators if any (best-effort)
        s = s.replace(",", "")
        v = float(s)
        if math.isnan(v) or math.isinf(v):
            return None
        return v
    except Exception:
        return None


def _fmt(x: Optional[float], decimals: int = 0) -> str:
    if x is None:
        return _MISSING_TEXT
    try:
        if math.isnan(float(x)) or math.isinf(float(x)):
            return _MISSING_TEXT
        # Use spaces for thousands separators to match your PPT style
        return f"{float(x):,.{decimals}f}".replace(",", " ")
    except Exception:
        return _MISSING_TEXT


def _pct_value_from_ratio(ratio_0_1: Optional[float], decimals: int = 0) -> str:
    """ratio_0_1 -> '12.3' (no % sign)."""
    if ratio_0_1 is None:
        return _MISSING_TEXT
    return _fmt(100.0 * ratio_0_1, decimals)


def _pct_value_from_pct(pct_value: Optional[float], decimals: int = 0) -> str:
    """pct_value already in percent units -> '12.3' (no % sign)."""
    if pct_value is None:
        return _MISSING_TEXT
    return _fmt(pct_value, decimals)


def _ensure_str(x: Any) -> str:
    if x is None:
        return ""
    return str(x)


def _pv_annuity(annual_cost: float, years: int, discount_rate: float, inflation_rate: float) -> float:
    """PV of annual cost stream with inflation (g) and discount (r), paid end of each year."""
    pv = 0.0
    for t in range(1, years + 1):
        pv += annual_cost * ((1 + inflation_rate) ** (t - 1)) / ((1 + discount_rate) ** t)
    return pv


def _classify_letter(kwh_m2_year: Optional[float]) -> str:
    """Simple EPC-like letter grade thresholds."""
    if kwh_m2_year is None:
        return _MISSING_TEXT
    v = float(kwh_m2_year)
    if v <= 50:
        return "A"
    if v <= 90:
        return "B"
    if v <= 130:
        return "C"
    if v <= 170:
        return "D"
    if v <= 210:
        return "E"
    if v <= 260:
        return "F"
    return "G"


# -----------------------------
# PPT shape utilities
# -----------------------------
def _iter_shapes(slide):
    """Iterate recursively over shapes including grouped shapes."""
    for shape in slide.shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            for subshape in shape.shapes:
                yield subshape


def _replace_tokens_in_text(text: str, token_map: Dict[str, str]) -> str:
    def repl(m):
        k = m.group(1).strip()
        return token_map.get(k, _MISSING_TEXT)

    return _TOKEN_PATTERN.sub(repl, text)


def _replace_tokens_in_shape(shape, token_map: Dict[str, str]) -> None:
    """
    Replace tokens even when split across runs (very common in PPTX).
    Also handles tables inside shapes.
    """
    if not getattr(shape, "has_text_frame", False):
        return

    tf = shape.text_frame
    if tf is None:
        return

    for p in tf.paragraphs:
        joined = "".join(r.text for r in p.runs)
        if "{{" in joined and "}}" in joined:
            replaced = _replace_tokens_in_text(joined, token_map)
            if len(p.runs) == 0:
                p.text = replaced
            else:
                p.runs[0].text = replaced
                for r in p.runs[1:]:
                    r.text = ""
        else:
            if "{{" in p.text and "}}" in p.text:
                p.text = _replace_tokens_in_text(p.text, token_map)

    if getattr(shape, "has_table", False):
        tbl = shape.table
        for row in tbl.rows:
            for cell in row.cells:
                if "{{" in cell.text and "}}" in cell.text:
                    cell.text = _replace_tokens_in_text(cell.text, token_map)


# -----------------------------
# Image placeholder detection helpers
# -----------------------------
def _normalize_placeholder_key(s: str) -> str:
    """
    ### FIX:
    Template visible text sometimes contains smart quotes or extra quotes:
      “IMG_RGB”, "IMG_RGB", ‘IMG_RGB’, etc.
    Normalize to IMG_RGB.
    """
    if not s:
        return ""
    t = s.strip()

    # remove common quote characters (straight + smart)
    t = t.strip('"\'')

    t = t.replace("“", "").replace("”", "")
    t = t.replace("‘", "").replace("’", "")
    t = t.replace("«", "").replace("»", "")
    t = t.replace("„", "").replace("‟", "")

    # collapse whitespace
    t = re.sub(r"\s+", " ", t).strip()

    return t


def _shape_alt_text(shape) -> str:
    """
    Alt-text (Description) used to identify image placeholders.

    ### FIX:
    Your previous XPath ".//p:cNvPr" fails without namespace bindings
    and gets swallowed by the try/except (returns "").

    Use namespace-agnostic XPath via local-name().
    """
    try:
        el = getattr(shape, "element", None) or getattr(shape, "_element", None)
        if el is None:
            return ""

        nodes = el.xpath(".//*[local-name()='cNvPr']")
        if not nodes:
            return ""

        descr = nodes[0].get("descr") or ""
        return _normalize_placeholder_key(descr)
    except Exception:
        return ""


def _shape_visible_text(shape) -> str:
    """Visible text inside the shape (used as fallback placeholder detection)."""
    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame and shape.text_frame.text:
            return _normalize_placeholder_key(shape.text_frame.text)
    except Exception:
        pass
    return ""


def _add_image_over_shape(slide, shape, image_bytes: bytes, remove_placeholder: bool = True) -> None:
    """
    Adds picture exactly over the placeholder shape.
    Optionally removes the placeholder so its 'IMG_*' text doesn't remain.
    """
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=width, height=height)

    if remove_placeholder:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception:
            # If removal fails, it's still fine (image should be on top).
            pass


def _b64_to_bytes(b64str: Optional[str]) -> Optional[bytes]:
    if not b64str:
        return None
    s = str(b64str).strip()
    if s.startswith("data:image"):
        s = s.split(",", 1)[-1]
    try:
        return base64.b64decode(s)
    except Exception:
        return None


# -----------------------------
# Label-based patching (optional)
# -----------------------------
def _norm_txt(s: Optional[str]) -> str:
    if not s:
        return ""
    return re.sub(r"\s+", " ", s.strip().lower())


def _set_value_right_of_label(slide, label_text: str, value_text: str, tol_y: int = 400000) -> bool:
    """
    Finds the shape whose *full* text equals label_text (case-insensitive),
    then sets the closest shape to its right (same row) to value_text.
    """
    label_text_n = _norm_txt(label_text)
    labels = []
    for sh in _iter_shapes(slide):
        if getattr(sh, "has_text_frame", False) and sh.text_frame and sh.text_frame.text:
            if _norm_txt(sh.text_frame.text) == label_text_n:
                labels.append(sh)

    if not labels:
        return False

    label = sorted(labels, key=lambda s: (s.top, s.left))[0]
    candidates = []
    for sh in _iter_shapes(slide):
        if sh is label:
            continue
        if not getattr(sh, "has_text_frame", False):
            continue
        if sh.left <= label.left:
            continue
        if abs(sh.top - label.top) > tol_y:
            continue
        candidates.append(sh)

    if not candidates:
        return False

    best = min(candidates, key=lambda s: (s.left - label.left))
    tf = best.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = value_text
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
    else:
        tf.text = value_text
    return True


# -----------------------------
# Slide 8: append => <LETTER> to kWh/m²·year lines and emphasize the letter
# -----------------------------
def _enhance_slide8_eec(prs: Presentation, token_map: Dict[str, str]) -> None:
    try:
        slide = prs.slides[7]  # slide 8 (0-indexed)
    except Exception:
        return

    facade_letter = token_map.get("EEC_LETTER_FACADE") or token_map.get("FACADE_EE_LETTER") or _MISSING_TEXT
    wall_letter = token_map.get("EEC_LETTER_WALL") or token_map.get("WALL_EE_LETTER") or _MISSING_TEXT
    win_letter = token_map.get("EEC_LETTER_WINDOWS") or token_map.get("WINDOWS_EE_LETTER") or _MISSING_TEXT

    def pick_letter(txt: str) -> str:
        t = (txt or "").lower()
        if "façade" in t or "facade" in t:
            return facade_letter
        if "window" in t or "glazing" in t:
            return win_letter
        if "wall" in t:
            return wall_letter
        return _MISSING_TEXT

    for shape in _iter_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        if tf is None:
            continue

        for p in tf.paragraphs:
            txt = p.text or ""
            low = txt.lower()

            if "kwh" not in low:
                continue
            if ("m²" not in txt) and ("m2" not in low):
                continue
            if "year" not in low:
                continue
            if "=>" in txt:
                continue

            letter = pick_letter(txt).strip()
            if not letter or letter == _MISSING_TEXT:
                continue

            base_size = None
            for r in p.runs:
                if r.font.size is not None:
                    base_size = r.font.size
                    break
            if base_size is None:
                base_size = Pt(16)

            r_arrow = p.add_run()
            r_arrow.text = "  =>  "
            r_arrow.font.size = base_size

            r_letter = p.add_run()
            r_letter.text = letter
            r_letter.font.bold = True
            r_letter.font.size = Pt(int(base_size.pt) + 6) if getattr(base_size, "pt", None) else Pt(22)


# -----------------------------
# Token map builder aligned to ThermalAI.pptx template tokens
# -----------------------------
def _build_token_map(report: Dict[str, Any], raw: Dict[str, Any]) -> Tuple[Dict[str, str], Dict[str, Any]]:
    """
    Returns (token_map, computed).
    """
    report = report or {}
    raw = raw or {}

    meta = report.get("meta") or {}
    inputs_r = report.get("inputs") or {}
    headline = report.get("headline") or {}

    raw_inputs = raw.get("inputs") or {}
    raw_results = raw.get("results") or {}
    raw_totals = raw_results.get("totals") or {}
    raw_comps = raw_results.get("components") or {}

    # ---- Basic identifiers
    analysis_id = _ensure_str(
        meta.get("analysis_id")
        or report.get("analysis_id")
        or raw.get("analysis_id")
        or meta.get("id")
        or ""
    )

    city = _ensure_str(meta.get("city") or raw_inputs.get("city") or "")
    country = _ensure_str(meta.get("country") or raw_inputs.get("country") or "")
    address = _ensure_str(meta.get("address") or raw_inputs.get("address") or "")
    google_maps = _ensure_str(meta.get("google_maps_link") or raw_inputs.get("google_maps_link") or "")

    # ---- Date/time (robust)
    dt_iso = (
        raw_inputs.get("datetime_iso")
        or (report.get("inputs") or {}).get("datetime_iso")
        or meta.get("datetime_iso")
        or raw_inputs.get("assessment_datetime")
        or meta.get("assessment_datetime")
    )

    date_of_report = (
        meta.get("date_of_report")
        or meta.get("report_date_iso")
        or (str(dt_iso)[:10] if dt_iso else None)
    )
    if not date_of_report:
        date_of_report = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    date_of_assessment = (
        meta.get("date_of_assessment")
        or meta.get("assessment_date")
        or (str(dt_iso)[:10] if dt_iso else None)
    )
    if not date_of_assessment:
        date_of_assessment = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    assessment_datetime = (
        meta.get("assessment_datetime")
        or dt_iso
        or datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    )

    # ---- Temperatures (°C)
    gps_lat = _safe_float(raw_inputs.get("latitude") or raw_inputs.get("gps_lat") or meta.get("latitude"))
    gps_lon = _safe_float(raw_inputs.get("longitude") or raw_inputs.get("gps_lon") or meta.get("longitude"))

    t_in = _safe_float(
        inputs_r.get("indoor_temp_c")
        or inputs_r.get("indoor_temperature_c")
        or raw_inputs.get("t_inside_c")
        or raw_inputs.get("t_inside")
    )
    if t_in is None:
        t_in = 22.0

    t_out = _safe_float(
        inputs_r.get("outdoor_temp_c")
        or inputs_r.get("outdoor_temperature_c")
        or raw_inputs.get("t_outside_c")
        or raw_inputs.get("t_outside")
    )

    # ΔT = Outdoor - Indoor
    dt = (t_out - t_in) if (t_out is not None and t_in is not None) else None

    # Optional climate stats
    t_out_min = _safe_float(raw_inputs.get("t_out_min_c") or raw_inputs.get("outdoor_temp_min_c"))
    t_out_max = _safe_float(raw_inputs.get("t_out_max_c") or raw_inputs.get("outdoor_temp_max_c"))
    t_out_avg = _safe_float(raw_inputs.get("t_out_avg_c") or raw_inputs.get("outdoor_temp_avg_c") or t_out)
    t_in_avg = _safe_float(raw_inputs.get("t_in_avg_c") or raw_inputs.get("indoor_temp_avg_c") or t_in)

    # ---- Areas
    facade_area_m2 = _safe_float(meta.get("facade_area_m2") or raw_inputs.get("facade_area_m2"))

    seg_counts = raw_inputs.get("segmentation_counts") or {}
    wall_px = float(_safe_float(seg_counts.get("wall_pixels")) or 0.0)
    win_px = float(_safe_float(seg_counts.get("window_pixels")) or 0.0)
    door_px = float(_safe_float(seg_counts.get("door_pixels")) or 0.0)
    total_px = wall_px + win_px + door_px

    wall_share = (wall_px / total_px) if total_px > 0 else None
    win_share = (win_px / total_px) if total_px > 0 else None

    wall_area_m2 = (facade_area_m2 * wall_share) if (facade_area_m2 is not None and wall_share is not None) else None
    win_area_m2 = (facade_area_m2 * win_share) if (facade_area_m2 is not None and win_share is not None) else None

    # ---- Degree hours -> HDD (K-day)
    degree_hours = _safe_float(raw_inputs.get("degree_hours_annual"))
    hdd_kday = (degree_hours / 24.0) if degree_hours is not None else _safe_float(meta.get("hdd") or raw_inputs.get("hdd"))

    # ---- Instantaneous W
    wall_inst_w = _safe_float(_dig(raw_comps, "wall.instantaneous_watts"))
    win_inst_w = _safe_float(_dig(raw_comps, "window.instantaneous_watts")) or _safe_float(_dig(raw_comps, "windows.instantaneous_watts"))
    door_inst_w = _safe_float(_dig(raw_comps, "door.instantaneous_watts")) or _safe_float(_dig(raw_comps, "doors.instantaneous_watts"))

    facade_inst_w = None
    if wall_inst_w is not None or win_inst_w is not None or door_inst_w is not None:
        facade_inst_w = (wall_inst_w or 0.0) + (win_inst_w or 0.0) + (door_inst_w or 0.0)

    # Heat flux (W/m²)
    facade_heat_flux = (facade_inst_w / facade_area_m2) if (facade_inst_w is not None and facade_area_m2) else None
    wall_heat_flux = (wall_inst_w / wall_area_m2) if (wall_inst_w is not None and wall_area_m2) else None
    window_heat_flux = (win_inst_w / win_area_m2) if (win_inst_w is not None and win_area_m2) else None

    # Heat loss coefficients as shares (0–1)
    if facade_inst_w is not None and facade_inst_w > 0:
        wall_coeff = (wall_inst_w / facade_inst_w) if wall_inst_w is not None else None
        win_coeff = (win_inst_w / facade_inst_w) if win_inst_w is not None else None
        facade_coeff = 1.0
    else:
        wall_coeff = win_coeff = facade_coeff = None

    # ---- Annual heat loss (kWh/year)
    wall_kwh = _safe_float(_dig(raw_comps, "wall.annual_kwh_delta")) or _safe_float(_dig(raw_comps, "wall.annual_kwh_u"))
    win_kwh = (
        _safe_float(_dig(raw_comps, "window.annual_kwh_delta"))
        or _safe_float(_dig(raw_comps, "windows.annual_kwh_delta"))
        or _safe_float(_dig(raw_comps, "window.annual_kwh_u"))
        or _safe_float(_dig(raw_comps, "windows.annual_kwh_u"))
    )

    annual_total_kwh = (
        _safe_float(headline.get("estimated_annual_heat_loss_kwh"))
        or _safe_float(raw_totals.get("annual_kwh_delta"))
        or _safe_float(raw_totals.get("annual_kwh_u"))
    )
    if annual_total_kwh is None and (wall_kwh is not None or win_kwh is not None):
        annual_total_kwh = (wall_kwh or 0.0) + (win_kwh or 0.0)

    facade_kwh_m2 = (annual_total_kwh / facade_area_m2) if (annual_total_kwh is not None and facade_area_m2) else None
    wall_kwh_m2 = (wall_kwh / wall_area_m2) if (wall_kwh is not None and wall_area_m2) else None
    win_kwh_m2 = (win_kwh / win_area_m2) if (win_kwh is not None and win_area_m2) else None

    # ---- CO2 (kg/y and kg/m²/y)
    emission_factor = _safe_float(
        inputs_r.get("emissions_factor_kg_per_kwh")
        or raw_inputs.get("emissions_factor_kg_per_kwh")
        or meta.get("emissions_factor_kg_per_kwh")
    )
    if emission_factor is None:
        emission_factor = 0.20

    annual_co2 = _safe_float(raw_totals.get("co2_kg_per_year")) or _safe_float(raw_totals.get("annual_co2_kg"))
    if annual_co2 is None and annual_total_kwh is not None:
        annual_co2 = annual_total_kwh * emission_factor

    annual_co2_m2 = (annual_co2 / facade_area_m2) if (annual_co2 is not None and facade_area_m2) else None
    annual_co2_kco2 = (annual_co2 / 1000.0) if annual_co2 is not None else None

    # ---- Economics
    energy_price = _safe_float(
        inputs_r.get("energy_price")
        or inputs_r.get("fuel_price_eur_per_kwh")
        or raw_inputs.get("fuel_price_eur_per_kwh")
        or meta.get("fuel_price_eur_per_kwh")
    )
    if energy_price is None:
        energy_price = 0.12

    discount_rate = _safe_float(inputs_r.get("discount_rate") or raw_inputs.get("discount_rate") or meta.get("discount_rate"))
    inflation_rate = _safe_float(inputs_r.get("inflation_rate") or raw_inputs.get("inflation_rate") or meta.get("inflation_rate"))
    if discount_rate is None:
        discount_rate = 0.03
    if inflation_rate is None:
        inflation_rate = 0.03

    wall_cost = (wall_kwh * energy_price) if (wall_kwh is not None and energy_price is not None) else None
    win_cost = (win_kwh * energy_price) if (win_kwh is not None and energy_price is not None) else None

    annual_cost_total = _safe_float(raw_totals.get("annual_cost_delta")) or _safe_float(raw_totals.get("annual_cost_u"))
    if annual_cost_total is None and (wall_cost is not None or win_cost is not None):
        annual_cost_total = (wall_cost or 0.0) + (win_cost or 0.0)
    if annual_cost_total is None and annual_total_kwh is not None:
        annual_cost_total = annual_total_kwh * energy_price

    pv_years = [1, 5, 10, 20, 30]
    pv_total: Dict[int, float] = {}
    pv_wall: Dict[int, float] = {}
    pv_win: Dict[int, float] = {}

    if annual_cost_total is not None:
        for y in pv_years:
            pv_total[y] = _pv_annuity(float(annual_cost_total), y, float(discount_rate), float(inflation_rate))
    if wall_cost is not None:
        for y in pv_years:
            pv_wall[y] = _pv_annuity(float(wall_cost), y, float(discount_rate), float(inflation_rate))
    if win_cost is not None:
        for y in pv_years:
            pv_win[y] = _pv_annuity(float(win_cost), y, float(discount_rate), float(inflation_rate))

    eec_letter_facade = _classify_letter(facade_kwh_m2)
    eec_letter_wall = _classify_letter(wall_kwh_m2)
    eec_letter_windows = _classify_letter(win_kwh_m2)

    tm: Dict[str, str] = {}

    def put(k: str, v: Any) -> None:
        s = _ensure_str(v).strip()
        tm[k] = s if s else _MISSING_TEXT

    put("analysis_id", analysis_id)
    put("ANALYSIS_ID", analysis_id)

    put("CITY", city)
    put("COUNTRY", country)

    put("ADDRESS_FULL", address)
    put("BUILDING_ADDRESS_FULL", address)
    put("ADDRESS", address)
    put("GOOGLE_MAPS_LINK", google_maps)

    put("DATE_OF_REPORT", date_of_report)
    put("REPORT_DATE", date_of_report)
    put("DATE_OF_ASSESSMENT", date_of_assessment)
    put("ASSESSMENT_DATETIME", assessment_datetime)

    put("INDOOR_TEMP_C", _fmt(t_in, 1))
    put("OUTDOOR_TEMP_C", _fmt(t_out, 1))
    put("OUTDOOR_TEMPERATURE_C", _fmt(t_out, 1))
    put("INDOOR_TEMPERATURE_C", _fmt(t_in, 1))

    put("T_OUT_MIN_C", _fmt(t_out_min, 1))
    put("T_OUT_MAX_C", _fmt(t_out_max, 1))
    put("T_OUT_AVG_C", _fmt(t_out_avg, 1))
    put("T_IN_AVG_C", _fmt(t_in_avg, 1))

    put("FACADE_SURFACE_TEMP_C", _fmt(_safe_float(raw_inputs.get("facade_surface_temp_c")), 1))

    put("GPS_LAT", _fmt(gps_lat, 6))
    put("GPS_LON", _fmt(gps_lon, 6))

    put("DELTA_T_C", _fmt(dt, 1))
    put("DELTA_T_K", _fmt(dt, 1))
    put("DT_C", _fmt(dt, 1))
    put("DT_K", _fmt(dt, 1))

    put("FACADE_AREA_M2", _fmt(facade_area_m2, 0))
    put("WALL_AREA_M2", _fmt(wall_area_m2, 0))
    put("WINDOW_AREA_M2", _fmt(win_area_m2, 0))
    put("WINDOWS_AREA_M2", _fmt(win_area_m2, 0))

    put("WINDOW_AREA_RATIO_PCT", _pct_value_from_ratio((win_area_m2 / facade_area_m2) if (win_area_m2 is not None and facade_area_m2) else None, 0))
    put("WALL_AREA_RATIO_PCT", _pct_value_from_ratio((wall_area_m2 / facade_area_m2) if (wall_area_m2 is not None and facade_area_m2) else None, 0))
    put("WINDOW_TO_FACADE_PCT", _pct_value_from_ratio(win_share, 0))
    put("WALL_TO_FACADE_PCT", _pct_value_from_ratio(wall_share, 0))

    put("HDD", _fmt(hdd_kday, 0))
    put("HEATING_DEGREE_DAYS_ANNUAL", _fmt(hdd_kday, 0))
    put("HDD_ANNUAL", _fmt(hdd_kday, 0))
    put("DEGREE_HOURS_ANNUAL", _fmt(degree_hours, 0))

    put("FACADE_HEAT_LOSS_W", _fmt(facade_inst_w, 0))
    put("INSTANTANEOUS_HEAT_LOSS_W", _fmt(facade_inst_w, 0))
    put("INSTANTANEOUS_HEAT_LOSS_MEASUREMENT_W", _fmt(facade_inst_w, 0))
    put("INSTANTANEOUS_HEAT_LOSS_MEASUREMENT_CONDITIONS_W", _fmt(facade_inst_w, 0))
    put("INSTANTANEOUS_HEAT_LOSS_MEASUREMENT_CONDITIONS", _fmt(facade_inst_w, 0))

    put("WALL_HEAT_LOSS_W", _fmt(wall_inst_w, 0))
    put("WINDOW_HEAT_LOSS_W", _fmt(win_inst_w, 0))

    put("FACADE_HEAT_FLUX_W_M2", _fmt(facade_heat_flux, 2))
    put("WALL_HEAT_FLUX_W_M2", _fmt(wall_heat_flux, 2))
    put("WINDOW_HEAT_FLUX_W_M2", _fmt(window_heat_flux, 2))

    # Heat loss share in % (template prints % sign)
    put("FACADE_HEAT_LOSS_COEFF", _pct_value_from_ratio(facade_coeff, 0))
    put("WALL_HEAT_LOSS_COEFF", _pct_value_from_ratio(wall_coeff, 0))
    put("WINDOW_HEAT_LOSS_COEFF", _pct_value_from_ratio(win_coeff, 0))

    put("FACADE_ANNUAL_HEAT_LOSS_KWH", _fmt(annual_total_kwh, 0))
    put("WINDOW_ANNUAL_HEAT_LOSS_KWH", _fmt(win_kwh, 0))
    put("WALL_ANNUAL_HEAT_LOSS_KWH", _fmt(wall_kwh, 0))

    put("FACADE_ANNUAL_HEAT_LOSS_KWH_M2", _fmt(facade_kwh_m2, 2))
    put("FACADE_ANNUAL_HEAT_LOSS_KWH_M2Y", _fmt(facade_kwh_m2, 2))
    put("WINDOW_ANNUAL_HEAT_LOSS_KWH_M2", _fmt(win_kwh_m2, 2))
    put("WINDOW_ANNUAL_HEAT_LOSS_KWH_M2Y", _fmt(win_kwh_m2, 2))
    put("WALL_ANNUAL_HEAT_LOSS_KWH_M2", _fmt(wall_kwh_m2, 2))
    put("WALL_ANNUAL_HEAT_LOSS_KWH_M2Y", _fmt(wall_kwh_m2, 2))

    put("ANNUAL_HEAT_LOSS_FACADE_KWH", _fmt(annual_total_kwh, 0))
    put("ANNUAL_HEAT_LOSS_WINDOWS_KWH", _fmt(win_kwh, 0))
    put("ANNUAL_HEAT_LOSS_WALL_KWH", _fmt(wall_kwh, 0))
    put("ANNUAL_HEAT_LOSS_FACADE_KWH_M2", _fmt(facade_kwh_m2, 2))
    put("ANNUAL_HEAT_LOSS_WINDOWS_KWH_M2", _fmt(win_kwh_m2, 2))
    put("ANNUAL_HEAT_LOSS_WALL_KWH_M2", _fmt(wall_kwh_m2, 2))

    put("CO2_FACADE_KG_PER_YEAR", _fmt(annual_co2, 0))
    put("CO2_FACADE_KG_PER_M2_YEAR", _fmt(annual_co2_m2, 3))
    put("FACADE_ANNUAL_CO2_KG", _fmt(annual_co2, 0))
    put("FACADE_ANNUAL_CO2_KG_M2", _fmt(annual_co2_m2, 3))
    put("FACADE_ANNUAL_CO2_KG_M2Y", _fmt(annual_co2_m2, 3))

    put("CO2_FACADE_KCO2_PER_YEAR", _fmt(annual_co2_kco2, 2))
    put("FACADE_ANNUAL_CO2_KCO2", _fmt(annual_co2_kco2, 2))

    put("EEC_LETTER_FACADE", eec_letter_facade)
    put("EEC_LETTER_WINDOWS", eec_letter_windows)
    put("EEC_LETTER_WALL", eec_letter_wall)
    put("EEC_KWH_M2_YEAR_FACADE", _fmt(facade_kwh_m2, 2))
    put("EEC_KWH_M2_YEAR_WINDOWS", _fmt(win_kwh_m2, 2))
    put("EEC_KWH_M2_YEAR_WALL", _fmt(wall_kwh_m2, 2))

    put("ENERGY_PRICE", _fmt(energy_price, 3))
    put("ENERGY_PRICE_CURRENT", _fmt(energy_price, 3))
    put("DISCOUNT_RATE_PCT", _pct_value_from_pct(discount_rate * 100.0, 2))
    put("ANNUAL_ENERGY_PRICE_INFLATION_PCT", _pct_value_from_pct(inflation_rate * 100.0, 2))
    put("ENERGY_PRICE_INFLATION_PCT", _pct_value_from_pct(inflation_rate * 100.0, 2))
    put("ENERGY_INFLATION_PCT", _pct_value_from_pct(inflation_rate * 100.0, 2))

    put("PV_TOTAL_1Y_EUR", _fmt(pv_total.get(1), 0))
    put("PV_TOTAL_Y1_EUR", _fmt(pv_total.get(1), 0))
    put("PV_TOTAL_5Y_EUR", _fmt(pv_total.get(5), 0))
    put("PV_TOTAL_10Y_EUR", _fmt(pv_total.get(10), 0))
    put("PV_TOTAL_20Y_EUR", _fmt(pv_total.get(20), 0))
    put("PV_TOTAL_30Y_EUR", _fmt(pv_total.get(30), 0))

    put("PV_WALL_1Y_EUR", _fmt(pv_wall.get(1), 0))
    put("PV_WALL_Y1_EUR", _fmt(pv_wall.get(1), 0))
    put("PV_WALL_5Y_EUR", _fmt(pv_wall.get(5), 0))
    put("PV_WALL_10Y_EUR", _fmt(pv_wall.get(10), 0))
    put("PV_WALL_20Y_EUR", _fmt(pv_wall.get(20), 0))
    put("PV_WALL_30Y_EUR", _fmt(pv_wall.get(30), 0))

    put("PV_WINDOWS_1Y_EUR", _fmt(pv_win.get(1), 0))
    put("PV_WINDOWS_Y1_EUR", _fmt(pv_win.get(1), 0))
    put("PV_WINDOWS_5Y_EUR", _fmt(pv_win.get(5), 0))
    put("PV_WINDOWS_10Y_EUR", _fmt(pv_win.get(10), 0))
    put("PV_WINDOWS_20Y_EUR", _fmt(pv_win.get(20), 0))
    put("PV_WINDOWS_30Y_EUR", _fmt(pv_win.get(30), 0))

    put("PV_WINDOW_1Y_EUR", _fmt(pv_win.get(1), 0))
    put("PV_WINDOW_5Y_EUR", _fmt(pv_win.get(5), 0))
    put("PV_WINDOW_10Y_EUR", _fmt(pv_win.get(10), 0))
    put("PV_WINDOW_20Y_EUR", _fmt(pv_win.get(20), 0))
    put("PV_WINDOW_30Y_EUR", _fmt(pv_win.get(30), 0))

    building_type = raw_inputs.get("building_type") or meta.get("building_type")
    building_year = raw_inputs.get("building_year") or meta.get("building_year")
    floor_area = raw_inputs.get("floor_area_m2") or meta.get("floor_area_m2")
    envelope_area = raw_inputs.get("envelope_area_m2") or meta.get("envelope_area_m2")
    num_stories = raw_inputs.get("num_stories") or meta.get("num_stories")
    heating_system = raw_inputs.get("heating_system") or meta.get("heating_system")
    climate_zone = raw_inputs.get("climate_zone") or meta.get("climate_zone")
    hdd_manual = raw_inputs.get("hdd") or meta.get("hdd")

    put("BUILDING_TYPE", _ensure_str(building_type))
    put("YEAR_OF_CONSTRUCTION", _ensure_str(building_year))
    put("BUILDING_YEAR", _ensure_str(building_year))
    put("TOTAL_FLOOR_AREA_M2", _fmt(_safe_float(floor_area), 0))
    put("ENVELOPE_AREA_M2", _fmt(_safe_float(envelope_area), 0))
    put("NUMBER_OF_STORIES", _ensure_str(num_stories))
    put("NUM_STORIES", _ensure_str(num_stories))
    put("PRIMARY_HEATING_SYSTEM", _ensure_str(heating_system))
    put("CLIMATE_ZONE", _ensure_str(climate_zone))
    put("HDD_MANUAL", _fmt(_safe_float(hdd_manual), 0))

    rh = _safe_float(raw_inputs.get("outdoor_rh_percent") or meta.get("outdoor_rh_percent"))
    wind = _safe_float(raw_inputs.get("wind_speed_mps") or meta.get("wind_speed_mps"))
    sky = raw_inputs.get("sky_conditions") or meta.get("sky_conditions")

    put("OUTDOOR_RH_PCT", _fmt(rh, 0))
    put("OUTDOOR_RH_PERCENT", _fmt(rh, 0))
    put("WIND_SPEED_MPS", _fmt(wind, 1))
    put("SKY_CONDITIONS", _ensure_str(sky))
    put("SURVEY_DATETIME", _ensure_str(raw_inputs.get("datetime_iso") or dt_iso))

    put("TOTAL_HEAT_LOSS_SENTENCE_KWH_PER_YEAR", _fmt(annual_total_kwh, 0))
    put("TOTAL_CO2_SENTENCE_KCO2_PER_YEAR", _fmt(annual_co2_kco2, 2))

    computed = {
        "dt": dt,
        "annual_total_kwh": annual_total_kwh,
        "annual_co2_kco2": annual_co2_kco2,
    }
    return tm, computed


# -----------------------------
# Populate images (FIXED/ROBUST)
# -----------------------------
def _populate_images(prs: Presentation, report_data: Dict[str, Any]) -> None:
    """
    Robustly inserts the 4 images into placeholder shapes identified by:
      - Alt Text Description: IMG_RGB / IMG_OVERLAY / IMG_RGB_BOXED / IMG_THERMAL_BOXED
      - OR visible placeholder text if alt text is missing.
    Also searches for base64 images in multiple possible locations in report_data.
    """
    report_data = report_data or {}
    report = (report_data.get("report") or {}) if isinstance(report_data, dict) else {}
    raw = (report_data.get("raw") or {}) if isinstance(report_data, dict) else {}

    raw_art = (raw.get("artifacts") or {}) if isinstance(raw, dict) else {}
    report_images = (report.get("images") or {}) if isinstance(report, dict) else {}
    report_art = (report.get("artifacts") or {}) if isinstance(report, dict) else {}

    # Also allow images at top-level (some clients do this)
    top_art = report_data.get("artifacts") or {}
    top_images = report_data.get("images") or {}

    def pick(*vals):
        for v in vals:
            if v:
                return v
        return None

    # Try many keys (covers current app_improved + older variants)
    rgb_b64 = pick(
        report_images.get("rgb_png_base64"),
        report_images.get("rgb_base64_png"),
        report_art.get("rgb_image_base64_png"),
        report_art.get("rgb_base64_png"),
        raw_art.get("rgb_image_base64_png"),
        raw_art.get("rgb_base64_png"),
        top_images.get("rgb_png_base64"),
        top_art.get("rgb_image_base64_png"),
    )
    overlay_b64 = pick(
        report_images.get("overlay_png_base64"),
        report_images.get("overlay_base64_png"),
        report_art.get("overlay_image_base64_png"),
        report_art.get("overlay_base64_png"),
        raw_art.get("overlay_image_base64_png"),
        raw_art.get("overlay_base64_png"),
        top_images.get("overlay_png_base64"),
        top_art.get("overlay_image_base64_png"),
    )
    rgb_box_b64 = pick(
        report_images.get("rgb_hotspot_boxes_png_base64"),
        report_images.get("rgb_boxed_png_base64"),
        report_art.get("rgb_hotspot_boxes_base64_png"),
        report_art.get("rgb_boxed_base64_png"),
        raw_art.get("rgb_hotspot_boxes_base64_png"),
        raw_art.get("rgb_boxed_base64_png"),
        top_images.get("rgb_hotspot_boxes_png_base64"),
        top_art.get("rgb_hotspot_boxes_base64_png"),
    )
    thr_box_b64 = pick(
        report_images.get("thermal_hotspot_boxes_png_base64"),
        report_images.get("thermal_boxed_png_base64"),
        report_art.get("thermal_hotspot_boxes_base64_png"),
        report_art.get("thermal_boxed_base64_png"),
        raw_art.get("thermal_hotspot_boxes_base64_png"),
        raw_art.get("thermal_boxed_base64_png"),
        top_images.get("thermal_hotspot_boxes_png_base64"),
        top_art.get("thermal_hotspot_boxes_base64_png"),
    )

    rgb_bytes = _b64_to_bytes(rgb_b64)
    overlay_bytes = _b64_to_bytes(overlay_b64)
    rgb_box_bytes = _b64_to_bytes(rgb_box_b64)
    thr_box_bytes = _b64_to_bytes(thr_box_b64)

    # Map placeholder id -> bytes
    img_map = {
        "IMG_RGB": rgb_bytes,
        "IMG_OVERLAY": overlay_bytes,
        "IMG_RGB_BOXED": rgb_box_bytes,
        "IMG_THERMAL_BOXED": thr_box_bytes,
    }

    # ### FIX:
    # Use normalized keys for lookup (IMG_RGB etc)
    img_map_norm = {k: v for k, v in img_map.items()}

    for slide in prs.slides:
        for shape in list(_iter_shapes(slide)):
            alt = _shape_alt_text(shape)
            vis = _shape_visible_text(shape)

            key = None
            if alt in img_map_norm:
                key = alt
            elif vis in img_map_norm:
                key = vis

            if not key:
                continue

            img_bytes = img_map_norm.get(key)
            if not img_bytes:
                continue

            _add_image_over_shape(slide, shape, img_bytes, remove_placeholder=True)


# -----------------------------
# PPT report build
# -----------------------------
def build_ppt_report(template_pptx_path: str, out_path: str, report_data: Dict[str, Any]) -> str:
    template_path = Path(template_pptx_path).resolve()
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))

    report = (report_data or {}).get("report") or {}
    raw = (report_data or {}).get("raw") or {}

    token_map, computed = _build_token_map(report, raw)

    # Replace tokens in all shapes (including groups)
    for slide in prs.slides:
        for shape in _iter_shapes(slide):
            _replace_tokens_in_shape(shape, token_map)

    # Slide 8 letter styling (after replacement)
    _enhance_slide8_eec(prs, token_map)

    # Images (robust)
    _populate_images(prs, report_data)

    # Optional label patching
    try:
        dt_val = computed.get("dt")
        if dt_val is not None:
            slide3 = prs.slides[2]  # slide 3 (0-indexed)
            _set_value_right_of_label(slide3, "Temperature Differential (ΔT)", _fmt(dt_val, 1))
    except Exception:
        pass

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def _convert_pptx_to_pdf(pptx_path: str, out_dir: str) -> str:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) not found in PATH; cannot convert PPTX to PDF.")

    out_dir_path = Path(out_dir)
    out_dir_path.mkdir(parents=True, exist_ok=True)

    cmd = [
        soffice,
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir_path.resolve()),
        str(Path(pptx_path).resolve()),
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {proc.stderr.strip() or proc.stdout.strip()}")

    pdf_path = out_dir_path / (Path(pptx_path).stem + ".pdf")
    if not pdf_path.exists():
        raise RuntimeError("LibreOffice conversion finished but PDF not found.")
    return str(pdf_path)


# -----------------------------
# Public API (must match app_improved.py call signature)
# -----------------------------
def build_reports(
    template_pptx_path: str,
    out_dir: str,
    analysis_id: str,
    report_data: Optional[Dict[str, Any]] = None,
    export_pdf: bool = False,
    report: Optional[Dict[str, Any]] = None,
    raw: Optional[Dict[str, Any]] = None,
    **kwargs,
) -> Tuple[str, Optional[str]]:
    """
    Build PPTX (always). Optionally build PDF.
    Returns (pptx_path, pdf_path_or_None).

    Accepts both:
      - build_reports(..., report_data={"report": report, "raw": raw}, ...)
      - build_reports(..., report=report, raw=raw, ...)
    """
    out_dir_path = Path(out_dir)
    out_dir_path.mkdir(parents=True, exist_ok=True)

    if report_data is None:
        report_data = {"report": report or {}, "raw": raw or {}}

    pptx_path = out_dir_path / f"ThermalAI_Report_{analysis_id}.pptx"
    build_ppt_report(template_pptx_path=template_pptx_path, out_path=str(pptx_path), report_data=report_data)

    pdf_path: Optional[str] = None
    if export_pdf:
        pdf_path = _convert_pptx_to_pdf(str(pptx_path), str(out_dir_path))

    return str(pptx_path), pdf_path


def build_report(*args, **kwargs):
    # Legacy alias used by some code paths
    return build_reports(*args, **kwargs)
